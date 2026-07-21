"""Render a DiagramSpec as native PowerPoint shapes.

Native shapes, not a rasterised image, because faculty edit these decks: a box
they can drag and relabel is worth more than a prettier PNG they cannot touch.
It also keeps a text layer for screen readers and search, and needs no Chromium
in the worker image.

Scope is deliberately small. A first version that draws 3-8 nodes in one
direction well is more useful than one that draws every graph badly, so
`capability()` reports DEGRADED for anything outside that and the caller falls
back to bullets rather than emitting spaghetti.
"""
from __future__ import annotations

import math
from dataclasses import dataclass

from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

from app.modules.m03_course_kit.presentation import text_metrics as tm
from app.modules.m03_course_kit.presentation.theme import Theme
from app.modules.m03_course_kit.schemas import DiagramSpec

FULL = "FULL"
DEGRADED = "DEGRADED"

# Above this, boxes shrink past legibility at 13.3" wide and edges start crossing.
MAX_NODES = 8
MAX_EDGE_LABEL_NODES = 6   # edge labels need gaps; tight chains have none

# Connection site indices on a rectangle, as PowerPoint numbers them.
_TOP, _LEFT, _BOTTOM, _RIGHT = 0, 1, 2, 3


@dataclass(frozen=True)
class _Box:
    """A placed node, in inches."""
    node_id: str
    label: str
    left: float
    top: float
    width: float
    height: float


def capability(spec: DiagramSpec | None) -> str:
    """Whether this diagram can be drawn well, badly, or not at all."""
    if spec is None or not spec.nodes:
        return DEGRADED
    if len(spec.nodes) > MAX_NODES:
        return DEGRADED
    ids = {n.id for n in spec.nodes}
    if len(ids) != len(spec.nodes):
        return DEGRADED                      # duplicate ids — edges are ambiguous
    if any(e.source not in ids or e.target not in ids for e in spec.edges):
        return DEGRADED                      # dangling edge — the model hallucinated a node
    return FULL


def _arrowhead(connector, theme: Theme) -> None:
    """Put a triangular head on a connector.

    python-pptx has no arrowhead API, so the line's <a:ln> gets an <a:tailEnd>
    element directly. Without this every edge is an ambiguous line and a
    flowchart stops saying which way the flow goes.
    """
    line = connector.line
    line.color.rgb = theme.palette.edge
    line.width = Pt(1.5)
    ln = line._get_or_add_ln()
    tail = ln.makeelement(qn("a:tailEnd"), {"type": "triangle", "w": "med", "len": "med"})
    ln.append(tail)


def _node_shape(slide, box: _Box, theme: Theme):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(box.left), Inches(box.top), Inches(box.width), Inches(box.height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = theme.palette.node_fill
    shape.line.color.rgb = theme.palette.node_border
    shape.line.width = Pt(1.25)
    shape.shadow.inherit = False

    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    # Size the label to the box rather than the box to the label: the layout has
    # already committed to a grid, and a shape that grows to fit its text is a
    # shape that overlaps its neighbour.
    pt = tm.fit_font_size(
        box.label, box.width, box.height, max_pt=theme.type.body_tight, min_pt=8.0, bold=True,
    )
    para = tf.paragraphs[0]
    para.alignment = PP_ALIGN.CENTER
    run = para.add_run()
    run.text = box.label
    run.font.size = Pt(pt)
    run.font.bold = True
    run.font.color.rgb = theme.palette.node_text
    return shape


def _edge_label(slide, text: str, cx: float, cy: float, theme: Theme) -> None:
    w = max(0.5, tm.text_width_in(text, theme.type.chip) + tm.TEXTBOX_H_INSET_IN + 0.06)
    h = tm.line_height_in(theme.type.chip) + 0.08
    box = slide.shapes.add_textbox(Inches(cx - w / 2), Inches(cy - h / 2), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = False
    para = tf.paragraphs[0]
    para.alignment = PP_ALIGN.CENTER
    run = para.add_run()
    run.text = text
    run.font.size = Pt(theme.type.chip)
    run.font.color.rgb = theme.palette.edge_label


def _connect(slide, shapes: dict, edge, boxes: dict, theme: Theme, *, vertical: bool) -> None:
    src, dst = shapes.get(edge.source), shapes.get(edge.target)
    if src is None or dst is None:
        return
    a, b = boxes[edge.source], boxes[edge.target]

    if vertical:
        begin_site, end_site = (_BOTTOM, _TOP) if a.top <= b.top else (_TOP, _BOTTOM)
    else:
        begin_site, end_site = (_RIGHT, _LEFT) if a.left <= b.left else (_LEFT, _RIGHT)

    # Straight for a neighbouring hop, elbow for a jump (a cycle's return edge).
    adjacent = abs((a.top - b.top) if vertical else (a.left - b.left)) < (
        (a.height if vertical else a.width) * 2.2
    )
    kind = MSO_CONNECTOR.STRAIGHT if adjacent else MSO_CONNECTOR.ELBOW
    conn = slide.shapes.add_connector(kind, Emu(0), Emu(0), Emu(0), Emu(0))
    conn.begin_connect(src, begin_site)
    conn.end_connect(dst, end_site)
    _arrowhead(conn, theme)

    if edge.label and len(boxes) <= MAX_EDGE_LABEL_NODES and adjacent:
        if vertical:
            # In the gap between the two rows, offset clear of the arrow itself.
            cx = a.left + a.width / 2 + 0.55
            cy = (a.top + a.height + b.top) / 2
        else:
            # ABOVE the row, not between the boxes: a label is routinely wider
            # than the horizontal gap, so centring it in that gap puts it on top
            # of the node it is meant to be labelling.
            cx = (a.left + a.width + b.left) / 2
            cy = a.top - 0.16
        _edge_label(slide, edge.label, cx, cy, theme)


# ---------------------------------------------------------------------------
# Layouts — each returns boxes in inches within (left, top, width, height)
# ---------------------------------------------------------------------------

def _layout_row(spec, region) -> tuple[list[_Box], bool]:
    """Left-to-right chain."""
    left, top, width, height = region
    n = len(spec.nodes)
    gap = 0.42
    bw = (width - gap * (n - 1)) / n
    bh = min(1.5, height * 0.55)
    ty = top + (height - bh) / 2
    return [
        _Box(node.id, node.label, left + i * (bw + gap), ty, bw, bh)
        for i, node in enumerate(spec.nodes)
    ], False


def _layout_column(spec, region) -> tuple[list[_Box], bool]:
    """Top-to-bottom chain, centred."""
    left, top, width, height = region
    n = len(spec.nodes)
    gap = 0.26
    bh = min(0.82, (height - gap * (n - 1)) / n)
    bw = min(6.2, width * 0.52)
    x = left + (width - bw) / 2
    total = n * bh + (n - 1) * gap
    ty = top + max(0.0, (height - total) / 2)
    return [
        _Box(node.id, node.label, x, ty + i * (bh + gap), bw, bh)
        for i, node in enumerate(spec.nodes)
    ], True


def _layout_cycle(spec, region) -> tuple[list[_Box], bool]:
    """Nodes on an ellipse, clockwise from the top."""
    left, top, width, height = region
    n = len(spec.nodes)
    bw, bh = min(2.7, width / 3.4), 0.78
    rx = (width - bw) / 2 * 0.92
    ry = (height - bh) / 2 * 0.92
    cx, cy = left + width / 2, top + height / 2
    boxes = []
    for i, node in enumerate(spec.nodes):
        angle = -math.pi / 2 + (2 * math.pi * i / n)
        boxes.append(_Box(
            node.id, node.label,
            cx + rx * math.cos(angle) - bw / 2,
            cy + ry * math.sin(angle) - bh / 2,
            bw, bh,
        ))
    return boxes, False


def _layout_hierarchy(spec, region) -> tuple[list[_Box], bool]:
    """Roots on top, children beneath, one row per depth."""
    left, top, width, height = region
    targets = {e.target for e in spec.edges}
    roots = [n.id for n in spec.nodes if n.id not in targets] or [spec.nodes[0].id]

    depth: dict[str, int] = {r: 0 for r in roots}
    children: dict[str, list[str]] = {}
    for e in spec.edges:
        children.setdefault(e.source, []).append(e.target)
    # Breadth-first: a node sits one row below its shallowest parent.
    queue = list(roots)
    while queue:
        nid = queue.pop(0)
        for child in children.get(nid, []):
            if child not in depth:
                depth[child] = depth[nid] + 1
                queue.append(child)
    for node in spec.nodes:              # unreachable nodes go on the bottom row
        depth.setdefault(node.id, max(depth.values(), default=0))

    rows: dict[int, list] = {}
    for node in spec.nodes:
        rows.setdefault(depth[node.id], []).append(node)

    n_rows = max(rows) + 1
    gap_y = 0.34
    bh = min(0.8, (height - gap_y * (n_rows - 1)) / n_rows)
    boxes = []
    for level, nodes in rows.items():
        gap_x = 0.34
        bw = min(3.0, (width - gap_x * (len(nodes) - 1)) / len(nodes))
        row_w = len(nodes) * bw + (len(nodes) - 1) * gap_x
        x0 = left + (width - row_w) / 2
        y = top + level * (bh + gap_y)
        for i, node in enumerate(nodes):
            boxes.append(_Box(node.id, node.label, x0 + i * (bw + gap_x), y, bw, bh))
    return boxes, True


def _layout_grid(spec, region) -> tuple[list[_Box], bool]:
    """Block diagram: a grid, up to 4 across."""
    left, top, width, height = region
    n = len(spec.nodes)
    cols = min(4, n) if n <= 4 else math.ceil(n / 2)
    rows = math.ceil(n / cols)
    gap = 0.34
    bw = (width - gap * (cols - 1)) / cols
    bh = min(1.35, (height - gap * (rows - 1)) / rows)
    boxes = []
    for i, node in enumerate(spec.nodes):
        r, c = divmod(i, cols)
        boxes.append(_Box(node.id, node.label,
                          left + c * (bw + gap), top + r * (bh + gap), bw, bh))
    return boxes, False


def _choose_layout(spec: DiagramSpec):
    dtype = (spec.diagram_type or "flowchart").lower()
    n = len(spec.nodes)
    if dtype == "cycle":
        return _layout_cycle
    if dtype == "hierarchy":
        return _layout_hierarchy
    if dtype == "block":
        return _layout_grid
    # flowchart / sequence: a short chain reads best across, a long one down.
    return _layout_row if n <= 4 else _layout_column


def render(slide, spec: DiagramSpec, region: tuple[float, float, float, float],
           theme: Theme) -> bool:
    """Draw `spec` into `region` (left, top, width, height in inches).

    Returns False if it declined to draw, so the caller can fall back to text
    rather than ship a slide with an empty hole in it.
    """
    if capability(spec) != FULL:
        return False
    try:
        layout = _choose_layout(spec)
        boxes, vertical = layout(spec, region)
        by_id = {b.node_id: b for b in boxes}
        shapes = {b.node_id: _node_shape(slide, b, theme) for b in boxes}
        for edge in spec.edges:
            _connect(slide, shapes, edge, by_id, theme, vertical=vertical)
        return True
    except Exception:
        # A diagram is an enhancement; it must never fail an export job.
        import logging
        logging.getLogger("vidya.m03.presentation.diagram").exception(
            "m03.diagram: render failed (type=%s nodes=%d) — falling back to text",
            spec.diagram_type, len(spec.nodes),
        )
        return False
