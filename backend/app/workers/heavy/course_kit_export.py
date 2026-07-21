"""
M03 course kit export task â€” PPTX (slide deck) / PDF (handout).

Runs on the celery-heavy queue.  Export is only permitted for
PUBLISHED and ARCHIVED kits; the service layer enforces this before
dispatching, and this task double-checks on entry.

Sensitive field gating (role-aware):
  - speaker_notes  â€” omitted for DEAN
  - model_answer   â€” omitted for DEAN

Assignments are the only assessment artifact this module exports.

Data loaded per export:
  - Kit detail (slides, assignments, JSONB plan fields)
  - Syllabus (for course_id)
  - Course from M01 (code, title)
"""
from __future__ import annotations

import asyncio
import io
import logging
import re
import sys
import uuid as uuid_module
from uuid import UUID

from app.database import tenant_schema_scope
from app.workers.base_task import VidyaTask
from app.workers.celery_app import celery_app

logger = logging.getLogger("vidya.worker.m03.export")

_async_engine = None


def _get_async_engine():
    global _async_engine
    if _async_engine is None:
        from sqlalchemy.ext.asyncio import create_async_engine
        from sqlalchemy.pool import NullPool
        from app.config import settings
        # NullPool: no connection caching between asyncio.run() calls.
        # On Windows --pool=solo each task runs in a fresh event loop; pooled
        # asyncpg connections attached to the previous (closed) loop raise
        # "Future attached to a different loop". NullPool prevents this.
        _async_engine = create_async_engine(settings.DATABASE_URL, poolclass=NullPool)
        # Re-apply the schema at the START OF EVERY TRANSACTION, not once per
        # session. A commit hands this connection back — NullPool closes it, a pool
        # recycles it — so anything after the first commit would otherwise run with
        # search_path = public, and a pooled connection could arrive still carrying
        # ANOTHER tenant's search_path. A commit cannot undo a per-BEGIN SET LOCAL.
        from app.database import bind_tenant_search_path
        bind_tenant_search_path(_async_engine)
    return _async_engine


# ---------------------------------------------------------------------------
# Celery task
# ---------------------------------------------------------------------------

@celery_app.task(
    base=VidyaTask,
    name="app.workers.heavy.export_course_kit",
    autoretry_for=(),
    max_retries=0,
)
def export_course_kit(
    *,
    job_id: str,
    kit_id: str,
    tenant_id: str,
    schema_name: str,
    export_format: str,
    requested_by_user_id: str,
    requested_by_role: str = "FACULTY",
    request_id: str | None = None,
    **kwargs,
) -> dict:
    if sys.platform == "win32":
        asyncio.set_event_loop_policy(asyncio.WindowsSelectorEventLoopPolicy())
    # Which tenant every transaction in this task belongs to. Held for the whole
    # run and dropped at the end of it: a worker process is long-lived and serves
    # every tenant in turn, and a schema left set is one the next task inherits.
    with tenant_schema_scope(schema_name):
        return asyncio.run(
            _run_export(
                kit_id=UUID(kit_id),
                tenant_id=UUID(tenant_id),
                schema_name=schema_name,
                export_format=export_format,
                requested_by_user_id=UUID(requested_by_user_id),
                requested_by_role=requested_by_role,
            )
        )


# ---------------------------------------------------------------------------
# Inner async implementation
# ---------------------------------------------------------------------------

async def _run_export(
    kit_id: UUID,
    tenant_id: UUID,
    schema_name: str,
    export_format: str,
    requested_by_user_id: UUID,
    requested_by_role: str,
) -> dict:
    from sqlalchemy import text
    from sqlalchemy.ext.asyncio import AsyncSession

    from app.core.audit_log.models import AuditEventType
    from app.core.audit_log.service import AuditService
    from app.core.storage.models import StorageEntityType
    from app.core.storage.repository import StorageRepository
    from app.modules.m01_program_advisor.repository import CourseRepository
    from app.modules.m02_syllabus.repository import SyllabusRepository
    from app.modules.m03_course_kit.models import CourseKitStatus
    from app.modules.m03_course_kit.repository import CourseKitRepository

    tenant_slug = schema_name.removeprefix("tenant_")
    engine = _get_async_engine()
    is_dean = requested_by_role == "DEAN"

    _EXPORTABLE = {CourseKitStatus.PUBLISHED, CourseKitStatus.ARCHIVED}

    try:
        async with AsyncSession(engine, expire_on_commit=False) as session:
            kit = await CourseKitRepository.get_detail(kit_id, db=session)
            if kit is None:
                raise ValueError(f"Course kit {kit_id} not found.")
            if kit.status not in _EXPORTABLE:
                raise ValueError(
                    f"Course kit {kit_id} is {kit.status.value}; "
                    "export requires PUBLISHED or ARCHIVED status."
                )

            syllabus = await SyllabusRepository.get_by_id(kit.syllabus_id, db=session)
            if syllabus is None:
                raise ValueError(f"Syllabus {kit.syllabus_id} not found.")

            course = await CourseRepository.get_by_id(syllabus.course_id, db=session)
            if course is None:
                raise ValueError(f"Course {syllabus.course_id} not found in M01.")

            # Fetch enrichment data for professional PPTX (tenant branding, COs, faculty name)
            _t_row = (await session.execute(
                text("SELECT name, logo_url, primary_color FROM public.tenants WHERE id = :tid"),
                {"tid": str(tenant_id)},
            )).mappings().one_or_none()
            _pptx_inst   = _t_row["name"]          if _t_row else None
            _pptx_logo   = _t_row["logo_url"]       if _t_row else None
            _pptx_color  = _t_row["primary_color"]  if _t_row else None

            from app.modules.m02_syllabus.repository import CourseOutcomeRepository
            _pptx_cos = await CourseOutcomeRepository.list_by_syllabus(syllabus.id, db=session)

            _u_row = (await session.execute(
                text("SELECT full_name FROM users WHERE id = :uid"),
                {"uid": str(requested_by_user_id)},
            )).mappings().one_or_none()
            _pptx_faculty = _u_row["full_name"] if _u_row else None

            buf = io.BytesIO()
            if export_format == "pptx":
                content_type = (
                    "application/vnd.openxmlformats-officedocument"
                    ".presentationml.presentation"
                )
                ext = "pptx"
                _generate_pptx(buf, kit, course, is_dean=is_dean,
                                cos=_pptx_cos, institution_name=_pptx_inst,
                                logo_url=_pptx_logo, primary_color=_pptx_color,
                                faculty_name=_pptx_faculty)
            elif export_format == "handout":
                content_type = "application/pdf"
                ext = "pdf"
                _generate_handout_pdf(buf, kit, course)
            else:  # "pdf"
                content_type = "application/pdf"
                ext = "pdf"
                _generate_pdf(buf, kit, course, is_dean=is_dean)

            buf.seek(0)
            file_bytes = buf.read()
            size_bytes = len(file_bytes)

            safe_code = re.sub(r"[^a-z0-9_-]", "_", course.code.lower())[:20].strip("_")
            suffix    = "_handout" if export_format == "handout" else ""
            filename  = f"kit_{safe_code}_u{kit.unit_number}_v{kit.version}{suffix}.{ext}"
            file_uuid  = uuid_module.uuid4()
            object_key = (
                f"vidya-assets/{tenant_slug}/course_kit_export"
                f"/{kit_id}/{file_uuid}-{filename}"
            )

            await asyncio.to_thread(_s3_put_object, object_key, file_bytes, content_type)

            asset = await StorageRepository.create(
                uploaded_by_user_id=requested_by_user_id,
                entity_type=StorageEntityType.COURSE_KIT_EXPORT.value,
                entity_id=kit_id,
                object_key=object_key,
                original_filename=filename,
                size_bytes=size_bytes,
                content_type=content_type,
                db=session,
            )
            await session.commit()

        download_url = await StorageRepository.generate_presigned_get_url(
            object_key=object_key,
            expires_in_seconds=86_400,
        )

        await AuditService.log(
            AuditEventType.COURSE_KIT_EXPORT_COMPLETED,
            actor_user_id=requested_by_user_id,
            actor_role="SYSTEM",
            tenant_id=tenant_id,
            schema_name=schema_name,
            target_entity="CourseKit",
            target_id=str(kit_id),
            metadata={
                "asset_id":   str(asset.id),
                "format":     export_format,
                "size_bytes": size_bytes,
            },
        )

        logger.info(
            "m03.export: %s export complete (kit=%s size=%d)",
            export_format, kit_id, size_bytes,
        )

        return {
            "download_url": download_url,
            "asset_id":     str(asset.id),
            "format":       export_format,
            "size_bytes":   size_bytes,
        }

    except Exception as exc:
        try:
            await AuditService.log(
                AuditEventType.COURSE_KIT_EXPORT_FAILED,
                actor_user_id=requested_by_user_id,
                actor_role="SYSTEM",
                tenant_id=tenant_id,
                schema_name=schema_name,
                target_entity="CourseKit",
                target_id=str(kit_id),
                metadata={"error": str(exc)[:500], "format": export_format},
            )
        except Exception:
            logger.exception("m03.export: failed to log COURSE_KIT_EXPORT_FAILED audit")
        raise


# ---------------------------------------------------------------------------
# S3 direct upload (worker has credentials â€” no presigned PUT needed)
# ---------------------------------------------------------------------------

def _s3_put_object(object_key: str, file_bytes: bytes, content_type: str) -> None:
    import boto3
    from app.config import settings

    client = boto3.client(
        "s3",
        endpoint_url=settings.S3_ENDPOINT,
        aws_access_key_id=settings.S3_ACCESS_KEY,
        aws_secret_access_key=settings.S3_SECRET_KEY,
        region_name=settings.S3_REGION,
        use_ssl=settings.S3_USE_SSL,
    )
    s3_key = object_key.removeprefix(f"{settings.S3_BUCKET}/")
    client.put_object(
        Bucket=settings.S3_BUCKET,
        Key=s3_key,
        Body=file_bytes,
        ContentType=content_type,
    )


# ---------------------------------------------------------------------------
# PPTX generation â€” university-quality lecture deck
# ---------------------------------------------------------------------------

def _generate_pptx(buf, kit, course, *, is_dean: bool,
                   cos=None, institution_name=None, logo_url=None,
                   primary_color=None, faculty_name=None) -> None:
    from datetime import date as _date
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN

    from app.modules.m03_course_kit.presentation import diagram as _diagram
    from app.modules.m03_course_kit.presentation import text_metrics as _metrics
    from app.modules.m03_course_kit.presentation.theme import Theme
    from app.modules.m03_course_kit.schemas import DiagramSpec

    # Palette and type scale now live in m03/presentation/theme.py. The short
    # aliases below are views onto the theme, not a second definition of it —
    # they keep the ~40 existing call sites unchanged.
    _THEME = Theme.for_tenant(primary_color)
    _PAL, _TYPE = _THEME.palette, _THEME.type

    # â”€â”€ Theme colours â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€
    _NAV  = _PAL.navy   # navy  â€“ cover bg / header bars
    _ACC  = _PAL.accent   # blue  â€“ accent / key-concepts
    _TEAL = _PAL.teal   # teal  â€“ examples / resources
    _GRN  = _PAL.green   # green â€“ summary / definitions
    _ORG  = _PAL.orange   # orange â€“ quiz slides
    _WHT  = _PAL.white
    _TXT  = _PAL.text   # dark slate
    _GRY  = _PAL.grey   # medium grey
    _LGRY = _PAL.light_grey   # light grey (table alternates)

    prs = Presentation()
    W = prs.slide_width  = Inches(13.33)
    H = prs.slide_height = Inches(7.5)
    BLANK = prs.slide_layouts[6]

    def _clear_placeholders(slide) -> None:
        for ph in list(slide.placeholders):
            sp = ph._element
            sp.getparent().remove(sp)

    # â”€â”€ Low-level helpers â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€

    def _bg(slide, c: RGBColor):
        f = slide.background.fill
        f.solid(); f.fore_color.rgb = c

    def _rect(slide, l, t, w, h, c: RGBColor):
        s = slide.shapes.add_shape(1, l, t, w, h)
        s.fill.solid(); s.fill.fore_color.rgb = c
        s.line.fill.background()
        return s

    def _txt(slide, text: str, l, t, w, h, *,
             sz=16, bold=False, italic=False,
             color: RGBColor = None, align=PP_ALIGN.LEFT):
        box = slide.shapes.add_textbox(l, t, w, h)
        tf  = box.text_frame
        tf.word_wrap = True
        p   = tf.paragraphs[0]
        p.alignment = align
        r   = p.add_run()
        r.text = text
        r.font.size = Pt(sz); r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color or _TXT
        return box

    def _wrapped_lines(text: str, width_in: float, sz: int) -> int:
        """How many lines `text` wraps to in a box `width_in` inches wide at
        `sz` pt, so callers can space stacked blocks far enough apart to avoid
        overlapping — PowerPoint text boxes don't clip overflow, they render
        straight past their nominal height.

        Measured against the real font metrics rather than estimated from an
        average character width. The old estimate was well-tuned for average
        prose but wrong at the tails: a line of wide glyphs came back one line
        short (i.e. overlapping), and narrow text over-reserved space.
        """
        return _metrics.wrapped_lines(text, width_in, sz)

    def _line_h(sz: int) -> float:
        """Height of one line at `sz` pt, in inches — read from the font."""
        return _metrics.line_height_in(sz)

    def _fit_line(text: str, width_in: float, sz: int, *, max_lines: int = 1) -> str:
        """Truncate `text` with an ellipsis so it occupies at most `max_lines`.

        For the footnote strip at the bottom of a slide, which is boxed in on
        two sides: the Bloom/CO footer sits at x=8.9", and the slide ends at
        7.5" down. Left unbounded, a long summary either runs under the footer
        or off the bottom edge, so it is cut rather than moved.
        """
        if _metrics.wrapped_lines(text, width_in, sz) <= max_lines:
            return text
        cut = len(text)
        while cut > 1 and _metrics.wrapped_lines(text[:cut] + '…', width_in, sz) > max_lines:
            cut -= 1
        return text[:cut].rstrip() + '…'

    def _header(slide, title: str, subtitle: str = '', color: RGBColor = None):
        hc = color or _NAV
        _rect(slide, 0, 0, W, Inches(1.02), hc)
        if subtitle:
            _txt(slide, subtitle, Inches(0.5), Inches(0.06),
                 Inches(12), Inches(0.3), sz=9, color=_GRY)
        _txt(slide, title, Inches(0.5), Inches(0.3), Inches(11.5), Inches(0.65),
             sz=24, bold=True, color=_WHT)

    def _tbl_hdr(table, headers, col_widths, hdr_color: RGBColor):
        for ci, (hd, cw) in enumerate(zip(headers, col_widths)):
            table.columns[ci].width = cw
            cell = table.cell(0, ci)
            cell.text = hd
            cell.fill.solid(); cell.fill.fore_color.rgb = hdr_color
            for p2 in cell.text_frame.paragraphs:
                for r2 in p2.runs:
                    r2.font.size = Pt(10); r2.font.bold = True
                    r2.font.color.rgb = _WHT

    def _tbl_row(table, ri, values, alt_color: RGBColor = None):
        for ci, val in enumerate(values):
            cell = table.cell(ri, ci)
            cell.text = str(val)
            if ri % 2 == 0 and alt_color:
                cell.fill.solid(); cell.fill.fore_color.rgb = alt_color
            for p2 in cell.text_frame.paragraphs:
                for r2 in p2.runs:
                    r2.font.size = Pt(10); r2.font.color.rgb = _TXT

    effective_cos = cos or []
    slides_sorted = sorted(kit.slides or [], key=lambda s: s.slide_number)
    tp = kit.teaching_plan or []

    # â”€â”€ 1. COVER SLIDE â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€
    sl = prs.slides.add_slide(BLANK)
    _clear_placeholders(sl)
    _bg(sl, _NAV)
    _rect(sl, 0, 0, Inches(0.22), H, _ACC)     # left accent strip

    logo_placed = False
    if logo_url:
        try:
            import urllib.request, io as _io
            with urllib.request.urlopen(logo_url, timeout=5) as resp:
                logo_bytes = resp.read()
            sl.shapes.add_picture(_io.BytesIO(logo_bytes),
                Inches(10.8), Inches(0.3), height=Inches(0.85))
            logo_placed = True
        except Exception:
            pass

    if institution_name:
        _txt(sl, institution_name.upper(),
             Inches(0.55), Inches(0.38),
             Inches(10 if logo_placed else 12.3), Inches(0.38),
             sz=11, bold=True, color=RGBColor(0x94, 0xA3, 0xB8))

    _txt(sl, course.code, Inches(0.55), Inches(1.7),
         Inches(12.3), Inches(0.55), sz=20, bold=True, color=_ACC)
    _txt(sl, course.title, Inches(0.55), Inches(2.3),
         Inches(12.3), Inches(1.05), sz=36, bold=True, color=_WHT)
    _txt(sl, f"Unit {kit.unit_number}", Inches(0.55), Inches(3.55),
         Inches(10), Inches(0.55), sz=20, color=RGBColor(0x94, 0xA3, 0xB8))

    _rect(sl, Inches(0.55), Inches(4.3), Inches(10), Inches(0.03), _ACC)

    meta_parts = []
    if faculty_name and not is_dean:
        meta_parts.append(faculty_name)
    meta_parts += [f"Version {kit.version}", kit.complexity_level.value,
                   _date.today().strftime('%B %Y')]
    _txt(sl, '  Â·  '.join(meta_parts), Inches(0.55), Inches(4.48),
         Inches(12.3), Inches(0.45), sz=12, color=RGBColor(0x94, 0xA3, 0xB8))

    # â”€â”€ 2. OBJECTIVES SLIDE â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€
    sl = prs.slides.add_slide(BLANK)
    _clear_placeholders(sl)
    _bg(sl, _WHT)
    _header(sl, f"Unit {kit.unit_number} â€” Objectives & Schedule",
            subtitle=f'{course.code}: {course.title}')

    # COs (left column)
    _rect(sl, Inches(0.4), Inches(1.1), Inches(6.1), Inches(0.25), _ACC)
    _txt(sl, 'COURSE OUTCOMES', Inches(0.4), Inches(1.1),
         Inches(6.1), Inches(0.25), sz=8, bold=True, color=_WHT, align=PP_ALIGN.CENTER)
    for i, co in enumerate(effective_cos[:8]):
        yt = Inches(1.42 + i * 0.62)
        _rect(sl, Inches(0.4), yt, Inches(0.85), Inches(0.44), _LGRY)
        _txt(sl, getattr(co, 'code', ''), Inches(0.4), yt,
             Inches(0.85), Inches(0.44), sz=10, bold=True, color=_ACC, align=PP_ALIGN.CENTER)
        _txt(sl, getattr(co, 'description', ''),
             Inches(1.35), yt, Inches(4.9), Inches(0.3), sz=11, color=_TXT)
        bl = getattr(co, 'bloom_level', '')
        blstr = bl.value if hasattr(bl, 'value') else str(bl)
        _txt(sl, blstr, Inches(1.35), yt + Inches(0.3),
             Inches(4.9), Inches(0.2), sz=8, italic=True, color=_GRY)

    # Weekly schedule table (right column)
    if tp:
        _rect(sl, Inches(6.8), Inches(1.1), Inches(6.1), Inches(0.25), _NAV)
        _txt(sl, 'WEEKLY SCHEDULE', Inches(6.8), Inches(1.1),
             Inches(6.1), Inches(0.25), sz=8, bold=True, color=_WHT, align=PP_ALIGN.CENTER)
        n = min(len(tp), 9)
        tbl = sl.shapes.add_table(
            n + 1, 3, Inches(6.8), Inches(1.4), Inches(6.1), Inches(n * 0.56 + 0.44)
        ).table
        _tbl_hdr(tbl, ('Week', 'Topic', 'Hrs'),
                 (Inches(0.85), Inches(4.05), Inches(1.2)), _NAV)
        for ri, wk in enumerate(tp[:n], 1):
            _tbl_row(tbl, ri, [
                str(wk.get('week', ri)),
                str(wk.get('topic', ''))[:55],
                str(wk.get('hours', 'â€”')),
            ], _LGRY)

    # -- Shared footer helper -------------------------------------------------------
    def _add_slide_footer(slide, ks, content, *, is_dean: bool) -> None:
        code_snippet = (content.get('code_snippet') or '').strip()
        summary_txt  = (content.get('student_summary') or '').strip()
        footer = []
        if ks.bloom_level:
            bl = ks.bloom_level.value if hasattr(ks.bloom_level, 'value') else ks.bloom_level
            footer.append(f'Bloom: {bl}')
        if ks.co_reference:
            footer.append(f'CO: {ks.co_reference}')
        if footer:
            _txt(slide, '   |   '.join(footer),
                 Inches(8.9), Inches(6.88), Inches(4.2), Inches(0.38), sz=9, color=_GRY)
        notes = []
        if not is_dean:
            tn = (content.get('teaching_notes') or '').strip()
            if tn:
                notes.append(f'TEACHING NOTES:\n{tn}')
            if ks.speaker_notes:
                notes.append(f'SPEAKER NOTES:\n{ks.speaker_notes}')
        if footer:
            notes.append('  |  '.join(footer))
        if code_snippet:
            notes.append(f'CODE:\n{code_snippet}')
        if summary_txt:
            notes.append(f'STUDENT SUMMARY: {summary_txt}')
        if notes:
            slide.notes_slide.notes_text_frame.text = '\n\n'.join(notes)

    # -- Renderer: OBJECTIVES ---------------------------------------------------------
    def _render_objectives(sl, ks, content, *, is_dean: bool) -> None:
        bullets      = [str(b) for b in (content.get('bullets') or []) if b]
        key_concepts = [str(k) for k in (content.get('key_concepts') or []) if k]
        summary_txt  = (content.get('student_summary') or '').strip()
        _header(sl, ks.title or f'Slide {ks.slide_number}',
                subtitle=f'{course.code}  \xb7  Unit {kit.unit_number}  \xb7  Slide {ks.slide_number}',
                color=_GRN)
        _rect(sl, Inches(11.15), Inches(0.08), Inches(1.95), Inches(0.27), _GRN)
        _txt(sl, 'OBJECTIVES', Inches(11.15), Inches(0.1), Inches(1.95), Inches(0.27),
             sz=8, bold=True, color=_WHT, align=PP_ALIGN.CENTER)
        _txt(sl, 'By the end of this section, students will be able to:',
             Inches(0.45), Inches(1.12), Inches(12.5), Inches(0.3),
             sz=11, italic=True, color=_GRY)
        items = bullets or key_concepts
        oy = Inches(1.52)
        for i, obj in enumerate(items[:8], 1):
            if oy > Inches(6.5):
                break
            obj_str = str(obj)
            lines = _wrapped_lines(obj_str, 11.7, _TYPE.body)
            row_h = max(0.62, 0.10 + lines * _line_h(_TYPE.body))
            _rect(sl, Inches(0.45), oy, Inches(0.46), Inches(0.48), _GRN)
            _txt(sl, str(i), Inches(0.45), oy, Inches(0.46), Inches(0.48),
                 sz=15, bold=True, color=_WHT, align=PP_ALIGN.CENTER)
            _txt(sl, obj_str, Inches(1.02), oy + Inches(0.06),
                 Inches(11.7), Inches(row_h), sz=14, color=_TXT)
            oy += Inches(row_h)
        if not items:
            _txt(sl, 'Objectives will be listed here.',
                 Inches(0.45), Inches(1.52), Inches(12.5), Inches(0.4), sz=14, color=_GRY)
        if summary_txt:
            _txt(sl, _fit_line(f'\u21b3 {summary_txt}', 8.2, _TYPE.footnote),
                 Inches(0.45), Inches(6.85),
                 Inches(8.2), Inches(0.4), sz=_TYPE.footnote, italic=True, color=_GRY)
        _add_slide_footer(sl, ks, content, is_dean=is_dean)

    # -- Renderer: WORKED_EXAMPLE -----------------------------------------------------
    def _render_worked_example(sl, ks, content, *, is_dean: bool) -> None:
        bullets      = [str(b) for b in (content.get('bullets') or []) if b]
        examples     = [str(e) for e in (content.get('examples') or []) if e]
        code_snippet = (content.get('code_snippet') or '').strip()
        summary_txt  = (content.get('student_summary') or '').strip()
        _CODE_BG  = RGBColor(0x1E, 0x29, 0x3B)
        _CODE_TXT = RGBColor(0xE2, 0xE8, 0xF0)
        _header(sl, ks.title or f'Slide {ks.slide_number}',
                subtitle=f'{course.code}  \xb7  Unit {kit.unit_number}  \xb7  Slide {ks.slide_number}',
                color=_TEAL)
        _rect(sl, Inches(11.15), Inches(0.08), Inches(1.95), Inches(0.27), _TEAL)
        _txt(sl, 'WORKED EXAMPLE', Inches(11.15), Inches(0.1), Inches(1.95), Inches(0.27),
             sz=7, bold=True, color=_WHT, align=PP_ALIGN.CENTER)
        problem = bullets[0] if bullets else (examples[0] if examples else '')
        if problem:
            _rect(sl, Inches(0.4), Inches(1.08), Inches(12.5), Inches(0.27), _TEAL)
            _txt(sl, 'PROBLEM', Inches(0.5), Inches(1.08), Inches(2), Inches(0.27),
                 sz=9, bold=True, color=_WHT)
            _rect(sl, Inches(0.4), Inches(1.37), Inches(12.5), Inches(0.72), _LGRY)
            _txt(sl, problem, Inches(0.55), Inches(1.4),
                 Inches(12.2), Inches(0.64), sz=13, color=_TXT)
        steps = (bullets[1:] if len(bullets) > 1 else []) or examples
        if steps:
            _rect(sl, Inches(0.4), Inches(2.18), Inches(12.5), Inches(0.27), _NAV)
            _txt(sl, 'SOLUTION STEPS', Inches(0.5), Inches(2.18), Inches(3), Inches(0.27),
                 sz=9, bold=True, color=_WHT)
            sy = Inches(2.52)
            step_ceiling = Inches(4.95 if code_snippet else 6.5)
            max_steps = 4 if code_snippet else 5
            for i, step in enumerate(steps[:max_steps], 1):
                if sy > step_ceiling:
                    break
                lines  = _wrapped_lines(str(step), 11.9, _TYPE.body_tight)
                step_h = max(0.5, 0.08 + lines * _line_h(_TYPE.body_tight))
                _rect(sl, Inches(0.4), sy, Inches(0.4), Inches(0.4), _ACC)
                _txt(sl, str(i), Inches(0.4), sy, Inches(0.4), Inches(0.4),
                     sz=13, bold=True, color=_WHT, align=PP_ALIGN.CENTER)
                _txt(sl, step, Inches(0.92), sy + Inches(0.04),
                     Inches(11.9), Inches(step_h), sz=13, color=_TXT)
                sy += Inches(step_h)
        if code_snippet:
            _rect(sl, Inches(0.4), Inches(5.1), Inches(12.5), Inches(0.27), _CODE_BG)
            _txt(sl, 'CODE', Inches(0.5), Inches(5.1), Inches(1), Inches(0.27),
                 sz=8, bold=True, color=_WHT)
            _rect(sl, Inches(0.4), Inches(5.39), Inches(12.5), Inches(0.95), _CODE_BG)
            for li, line in enumerate(code_snippet.split('\n')[:3]):
                _txt(sl, line, Inches(0.55), Inches(5.43 + li * 0.28),
                     Inches(12.2), Inches(0.26), sz=10, italic=True, color=_CODE_TXT)
        if summary_txt:
            _txt(sl, _fit_line(f'\u21b3 {summary_txt}', 8.2, _TYPE.footnote),
                 Inches(0.45), Inches(6.85),
                 Inches(8.2), Inches(0.4), sz=_TYPE.footnote, italic=True, color=_GRY)
        _add_slide_footer(sl, ks, content, is_dean=is_dean)

    # -- Renderer: CODE ---------------------------------------------------------------
    def _render_code(sl, ks, content, *, is_dean: bool) -> None:
        code_snippet = (content.get('code_snippet') or '').strip()
        bullets      = [str(b) for b in (content.get('bullets') or []) if b]
        key_concepts = [str(k) for k in (content.get('key_concepts') or []) if k]
        summary_txt  = (content.get('student_summary') or '').strip()
        _CODE_BG  = RGBColor(0x1E, 0x29, 0x3B)
        _GUTTER   = RGBColor(0x0F, 0x17, 0x2A)
        _CODE_TXT = RGBColor(0xE2, 0xE8, 0xF0)
        _LINENO   = RGBColor(0x64, 0x74, 0x8B)
        _header(sl, ks.title or f'Slide {ks.slide_number}',
                subtitle=f'{course.code}  \xb7  Unit {kit.unit_number}  \xb7  Slide {ks.slide_number}',
                color=_NAV)
        _rect(sl, Inches(11.15), Inches(0.08), Inches(1.95), Inches(0.27), _ACC)
        _txt(sl, 'CODE', Inches(11.15), Inches(0.1), Inches(1.95), Inches(0.27),
             sz=8, bold=True, color=_WHT, align=PP_ALIGN.CENTER)
        _rect(sl, Inches(0.4), Inches(1.08), Inches(8.1), Inches(5.8), _CODE_BG)
        _rect(sl, Inches(0.4), Inches(1.08), Inches(0.52), Inches(5.8), _GUTTER)
        if code_snippet:
            for li, line in enumerate(code_snippet.split('\n')[:19]):
                ly = Inches(1.14 + li * 0.28)
                _txt(sl, str(li + 1), Inches(0.4), ly, Inches(0.52), Inches(0.27),
                     sz=8, color=_LINENO, align=PP_ALIGN.RIGHT)
                _txt(sl, line or ' ', Inches(0.97), ly, Inches(7.38), Inches(0.27),
                     sz=10, italic=True, color=_CODE_TXT)
        else:
            _txt(sl, '# No code snippet provided', Inches(0.97), Inches(1.2),
                 Inches(7.38), Inches(0.3), sz=10, italic=True, color=_LINENO)
        rx, ry = Inches(8.75), Inches(1.08)
        if bullets:
            _rect(sl, rx, ry, Inches(4.1), Inches(0.27), _ACC)
            _txt(sl, 'EXPLANATION', rx, ry, Inches(4.1), Inches(0.27),
                 sz=8, bold=True, color=_WHT, align=PP_ALIGN.CENTER)
            ry += Inches(0.34)
            for b in bullets[:6]:
                _txt(sl, f'\u2022 {b}', rx, ry, Inches(4.0), Inches(0.48),
                     sz=11, color=_TXT)
                ry += Inches(0.52)
            ry += Inches(0.1)
        if key_concepts:
            _rect(sl, rx, ry, Inches(4.1), Inches(0.27), _GRN)
            _txt(sl, 'CONCEPTS', rx, ry, Inches(4.1), Inches(0.27),
                 sz=8, bold=True, color=_WHT, align=PP_ALIGN.CENTER)
            ry += Inches(0.34)
            for kc in key_concepts[:4]:
                _txt(sl, f'\u2022 {kc}', rx, ry, Inches(4.0), Inches(0.38),
                     sz=11, bold=True, color=_ACC)
                ry += Inches(0.42)
        if summary_txt:
            _txt(sl, _fit_line(f'\u21b3 {summary_txt}', 8.2, _TYPE.footnote),
                 Inches(0.45), Inches(6.85),
                 Inches(8.2), Inches(0.4), sz=_TYPE.footnote, italic=True, color=_GRY)
        _add_slide_footer(sl, ks, content, is_dean=is_dean)

    # -- Renderer: COMMON_MISTAKES ----------------------------------------------------
    def _render_common_mistakes(sl, ks, content, *, is_dean: bool) -> None:
        bullets      = [str(b) for b in (content.get('bullets') or []) if b]
        key_concepts = [str(k) for k in (content.get('key_concepts') or []) if k]
        definitions  = [str(d) for d in (content.get('definitions') or []) if d]
        summary_txt  = (content.get('student_summary') or '').strip()
        _WRONG    = RGBColor(0x99, 0x1B, 0x1B)
        _RIGHT    = RGBColor(0x06, 0x55, 0x35)
        _WRONG_BG = RGBColor(0xFE, 0xF2, 0xF2)
        _RIGHT_BG = RGBColor(0xF0, 0xFD, 0xF4)
        _header(sl, ks.title or f'Slide {ks.slide_number}',
                subtitle=f'{course.code}  \xb7  Unit {kit.unit_number}  \xb7  Slide {ks.slide_number}',
                color=_ORG)
        _rect(sl, Inches(11.15), Inches(0.08), Inches(1.95), Inches(0.27), _ORG)
        _txt(sl, 'COMMON MISTAKES', Inches(11.15), Inches(0.1), Inches(1.95), Inches(0.27),
             sz=7, bold=True, color=_WHT, align=PP_ALIGN.CENTER)
        _rect(sl, Inches(0.4), Inches(1.08), Inches(6.05), Inches(0.3), _WRONG)
        _txt(sl, 'COMMON MISTAKE', Inches(0.5), Inches(1.08),
             Inches(5.9), Inches(0.3), sz=10, bold=True, color=_WHT)
        _rect(sl, Inches(6.75), Inches(1.08), Inches(6.05), Inches(0.3), _RIGHT)
        _txt(sl, 'CORRECT APPROACH', Inches(6.85), Inches(1.08),
             Inches(5.9), Inches(0.3), sz=10, bold=True, color=_WHT)
        if key_concepts:
            mistakes    = bullets
            corrections = key_concepts
        elif definitions:
            mistakes    = bullets
            corrections = definitions
        elif len(bullets) >= 2:
            mistakes    = bullets[::2]
            corrections = bullets[1::2]
        else:
            mistakes    = bullets
            corrections = []
        n_rows = min(max(len(mistakes), len(corrections), 1), 5)
        for i in range(n_rows):
            ry      = Inches(1.48 + i * 1.05)
            mistake = mistakes[i]    if i < len(mistakes)    else ''
            correct = corrections[i] if i < len(corrections) else ''
            if mistake:
                _rect(sl, Inches(0.4), ry, Inches(6.05), Inches(0.9), _WRONG_BG)
                _txt(sl, mistake, Inches(0.55), ry + Inches(0.1),
                     Inches(5.75), Inches(0.75), sz=12, color=_WRONG)
            if correct:
                _rect(sl, Inches(6.75), ry, Inches(6.05), Inches(0.9), _RIGHT_BG)
                _txt(sl, correct, Inches(6.9), ry + Inches(0.1),
                     Inches(5.75), Inches(0.75), sz=12, color=_RIGHT)
        if summary_txt:
            _txt(sl, _fit_line(f'\u21b3 {summary_txt}', 8.2, _TYPE.footnote),
                 Inches(0.45), Inches(6.85),
                 Inches(8.2), Inches(0.4), sz=_TYPE.footnote, italic=True, color=_GRY)
        _add_slide_footer(sl, ks, content, is_dean=is_dean)

    # -- Renderer: ACTIVITY -----------------------------------------------------------
    def _render_activity(sl, ks, content, *, is_dean: bool) -> None:
        activity    = (content.get('classroom_activity') or '').strip()
        bullets     = [str(b) for b in (content.get('bullets') or []) if b]
        definitions = [str(d) for d in (content.get('definitions') or []) if d]
        summary_txt = (content.get('student_summary') or '').strip()
        _ACT_BG = RGBColor(0xEC, 0xFD, 0xFF)
        _header(sl, ks.title or f'Slide {ks.slide_number}',
                subtitle=f'{course.code}  \xb7  Unit {kit.unit_number}  \xb7  Slide {ks.slide_number}',
                color=_TEAL)
        _rect(sl, Inches(11.15), Inches(0.08), Inches(1.95), Inches(0.27), _TEAL)
        _txt(sl, 'ACTIVITY', Inches(11.15), Inches(0.1), Inches(1.95), Inches(0.27),
             sz=8, bold=True, color=_WHT, align=PP_ALIGN.CENTER)
        top_y = Inches(1.08)
        if activity:
            _rect(sl, Inches(0.4), top_y, Inches(12.5), Inches(0.27), _TEAL)
            _txt(sl, 'ACTIVITY', Inches(0.5), top_y, Inches(2), Inches(0.27),
                 sz=9, bold=True, color=_WHT)
            _rect(sl, Inches(0.4), top_y + Inches(0.29), Inches(12.5), Inches(0.78), _ACT_BG)
            _txt(sl, activity, Inches(0.55), top_y + Inches(0.32),
                 Inches(12.2), Inches(0.7), sz=13, color=_TXT)
            top_y += Inches(1.17)
        steps = bullets
        if steps:
            _rect(sl, Inches(0.4), top_y, Inches(12.5), Inches(0.27), _NAV)
            _txt(sl, 'INSTRUCTIONS', Inches(0.5), top_y, Inches(3), Inches(0.27),
                 sz=9, bold=True, color=_WHT)
            sy = top_y + Inches(0.38)
            for i, step in enumerate(steps[:6], 1):
                _rect(sl, Inches(0.4), sy, Inches(0.38), Inches(0.38), _TEAL)
                _txt(sl, str(i), Inches(0.4), sy, Inches(0.38), Inches(0.38),
                     sz=12, bold=True, color=_WHT, align=PP_ALIGN.CENTER)
                _txt(sl, step, Inches(0.88), sy + Inches(0.03),
                     Inches(11.95), Inches(0.36), sz=13, color=_TXT)
                sy += Inches(0.52)
        elif not activity:
            _txt(sl, '\u2022 Complete the assigned activity',
                 Inches(0.5), Inches(1.5), Inches(12), Inches(0.4), sz=14, color=_TXT)
        if definitions:
            _rect(sl, Inches(0.4), Inches(5.7), Inches(12.5), Inches(0.27), _ACC)
            _txt(sl, 'DISCUSSION PROMPT', Inches(0.5), Inches(5.7),
                 Inches(4), Inches(0.27), sz=9, bold=True, color=_WHT)
            dy = Inches(6.02)
            for d in definitions[:2]:
                _txt(sl, f'\u2192 {d}', Inches(0.55), dy,
                     Inches(12.2), Inches(0.38), sz=12, italic=True, color=_TXT)
                dy += Inches(0.42)
        if summary_txt:
            _txt(sl, _fit_line(f'\u21b3 {summary_txt}', 8.2, _TYPE.footnote),
                 Inches(0.45), Inches(6.85),
                 Inches(8.2), Inches(0.4), sz=_TYPE.footnote, italic=True, color=_GRY)
        _add_slide_footer(sl, ks, content, is_dean=is_dean)

    # -- Renderer: QUIZ ---------------------------------------------------------------
    def _render_quiz(sl, ks, content, *, is_dean: bool) -> None:
        bullets     = [str(b) for b in (content.get('bullets') or []) if b]
        examples    = [str(e) for e in (content.get('examples') or []) if e]
        definitions = [str(d) for d in (content.get('definitions') or []) if d]
        summary_txt = (content.get('student_summary') or '').strip()
        _OPT_COLORS = [
            RGBColor(0xEF, 0xF6, 0xFF),
            RGBColor(0xF5, 0xF3, 0xFF),
            RGBColor(0xF0, 0xFD, 0xF4),
            RGBColor(0xFF, 0xF7, 0xED),
        ]
        _LBL_COLORS = [_ACC, RGBColor(0x74, 0x4E, 0xB8), _GRN, _ORG]
        _header(sl, ks.title or f'Slide {ks.slide_number}',
                subtitle=f'{course.code}  \xb7  Unit {kit.unit_number}  \xb7  Slide {ks.slide_number}',
                color=_ORG)
        _rect(sl, Inches(11.15), Inches(0.08), Inches(1.95), Inches(0.27), _ORG)
        _txt(sl, 'QUIZ', Inches(11.15), Inches(0.1), Inches(1.95), Inches(0.27),
             sz=8, bold=True, color=_WHT, align=PP_ALIGN.CENTER)
        question = bullets[0] if bullets else (examples[0] if examples else '')
        if question:
            _rect(sl, Inches(0.4), Inches(1.08), Inches(12.5), Inches(0.27), _ORG)
            _txt(sl, 'QUESTION', Inches(0.5), Inches(1.08),
                 Inches(2), Inches(0.27), sz=9, bold=True, color=_WHT)
            _rect(sl, Inches(0.4), Inches(1.37), Inches(12.5), Inches(0.88),
                  RGBColor(0xFF, 0xF7, 0xED))
            _txt(sl, question, Inches(0.55), Inches(1.42),
                 Inches(12.2), Inches(0.78), sz=15, bold=True, color=_TXT)
        options = (bullets[1:] if len(bullets) > 1 else []) or \
                  (examples[1:] if len(examples) > 1 else examples)
        for i, opt in enumerate(options[:4]):
            col = i % 2
            row = i // 2
            ox  = Inches(0.4)  if col == 0 else Inches(6.8)
            oy  = Inches(2.38) + Inches(row * 1.35)
            _rect(sl, ox, oy, Inches(6.1), Inches(1.1), _OPT_COLORS[i])
            _rect(sl, ox, oy, Inches(0.5), Inches(1.1), _LBL_COLORS[i])
            _txt(sl, ('A', 'B', 'C', 'D')[i], ox, oy, Inches(0.5), Inches(1.1),
                 sz=20, bold=True, color=_WHT, align=PP_ALIGN.CENTER)
            _txt(sl, opt, ox + Inches(0.58), oy + Inches(0.28),
                 Inches(5.38), Inches(0.6), sz=13, color=_TXT)
        answer_text = definitions[0] if definitions else (summary_txt or '')
        if not is_dean and answer_text:
            _rect(sl, Inches(0.4), Inches(6.38), Inches(12.5), Inches(0.3), _GRN)
            _txt(sl, f'ANSWER: {answer_text}', Inches(0.5), Inches(6.38),
                 Inches(12.2), Inches(0.3), sz=10, bold=True, color=_WHT)
        _add_slide_footer(sl, ks, content, is_dean=is_dean)

    # -- Renderer: SUMMARY content slide ----------------------------------------------
    def _render_summary_slide(sl, ks, content, *, is_dean: bool) -> None:
        bullets      = [str(b) for b in (content.get('bullets') or []) if b]
        key_concepts = [str(k) for k in (content.get('key_concepts') or []) if k]
        summary_txt  = (content.get('student_summary') or '').strip()
        _header(sl, ks.title or f'Slide {ks.slide_number}',
                subtitle=f'{course.code}  \xb7  Unit {kit.unit_number}  \xb7  Slide {ks.slide_number}',
                color=_GRN)
        _rect(sl, Inches(11.15), Inches(0.08), Inches(1.95), Inches(0.27), _GRN)
        _txt(sl, 'SUMMARY', Inches(11.15), Inches(0.1), Inches(1.95), Inches(0.27),
             sz=8, bold=True, color=_WHT, align=PP_ALIGN.CENTER)
        _txt(sl, 'Key takeaways from this section:',
             Inches(0.45), Inches(1.12), Inches(12.5), Inches(0.3),
             sz=11, italic=True, color=_GRY)
        items = bullets or key_concepts or ([summary_txt] if summary_txt else [])
        oy = Inches(1.52)
        for item in items[:8]:
            _rect(sl, Inches(0.45), oy + Inches(0.08), Inches(0.34), Inches(0.34), _GRN)
            _txt(sl, 'v', Inches(0.45), oy + Inches(0.06), Inches(0.34), Inches(0.34),
                 sz=12, bold=True, color=_WHT, align=PP_ALIGN.CENTER)
            _txt(sl, str(item), Inches(0.9), oy + Inches(0.08),
                 Inches(11.7), Inches(0.38), sz=14, color=_TXT)
            oy += Inches(0.62)
        if summary_txt and items and summary_txt not in items:
            _txt(sl, _fit_line(f'\u21b3 {summary_txt}', 8.2, _TYPE.footnote),
                 Inches(0.45), Inches(6.85),
                 Inches(8.2), Inches(0.4), sz=_TYPE.footnote, italic=True, color=_GRY)
        _add_slide_footer(sl, ks, content, is_dean=is_dean)

    # -- Renderer: DIAGRAM ------------------------------------------------------------
    def _parse_diagram(content: dict):
        """The slide's diagram, or None if there isn't a drawable one.

        Kits generated before DiagramSpec existed carry `diagram_prompt` — prose
        describing a picture, which cannot be drawn. Those still fall through to
        the text renderers, exactly as they do today.
        """
        raw = content.get('diagram')
        if not isinstance(raw, dict):
            return None
        try:
            spec = DiagramSpec.model_validate(raw)
        except Exception:
            return None
        return spec if _diagram.capability(spec) == _diagram.FULL else None

    def _render_diagram(sl, ks, content, *, is_dean: bool) -> None:
        """Diagram left, explanation right.

        The diagram is drawn as native PowerPoint shapes, so faculty can drag a
        box or fix a label. If it declines to draw, this slide falls back to the
        generic renderer rather than shipping a slide with a hole in it.
        """
        # Decide BEFORE drawing anything: the generic renderer paints its own
        # header, so handing over after this one is drawn stacks two of them.
        spec = _parse_diagram(content)
        if spec is None:
            _render_generic(sl, ks, content, is_dean=is_dean)
            return

        bullets      = [str(b) for b in (content.get('bullets') or []) if b]
        key_concepts = [str(k) for k in (content.get('key_concepts') or []) if k]
        summary_txt  = (content.get('student_summary') or '').strip()

        _header(sl, ks.title or f'Slide {ks.slide_number}',
                subtitle=f'{course.code}  \xb7  Unit {kit.unit_number}  \xb7  Slide {ks.slide_number}',
                color=_NAV)
        _rect(sl, Inches(11.15), Inches(0.08), Inches(1.95), Inches(0.27), _ACC)
        _txt(sl, 'DIAGRAM', Inches(11.15), Inches(0.1), Inches(1.95), Inches(0.27),
             sz=_TYPE.chip, bold=True, color=_WHT, align=PP_ALIGN.CENTER)

        has_side = bool(bullets or key_concepts)
        region = (0.45, 1.15, 8.3 if has_side else 12.45, 5.5)
        if not _diagram.render(sl, spec, region, _THEME):
            # capability() approved it but drawing raised. The header is already
            # down, so fill the space with the walkthrough rather than hand over
            # to a renderer that would draw a second header over this one.
            fy = Inches(1.3)
            for b in bullets[:7]:
                if fy > Inches(6.4):
                    break
                h = 0.06 + _wrapped_lines(f'• {b}', 12.4, _TYPE.body) * _line_h(_TYPE.body)
                _txt(sl, f'• {b}', Inches(0.5), fy, Inches(12.4), Inches(h),
                     sz=_TYPE.body, color=_TXT)
                fy += Inches(h + 0.06)
            _add_slide_footer(sl, ks, content, is_dean=is_dean)
            return

        ry = Inches(1.15)
        if bullets:
            _rect(sl, Inches(8.9), ry, Inches(4.1), Inches(0.25), _ACC)
            _txt(sl, 'HOW IT WORKS', Inches(8.9), ry, Inches(4.1), Inches(0.25),
                 sz=_TYPE.chip, bold=True, color=_WHT, align=PP_ALIGN.CENTER)
            ry += Inches(0.32)
            for b in bullets[:6]:
                if ry > Inches(5.9):
                    break
                lines = _wrapped_lines(f'• {b}', 4.0, _TYPE.caption)
                h = 0.06 + lines * _line_h(_TYPE.caption)
                _txt(sl, f'• {b}', Inches(8.95), ry, Inches(4.0), Inches(h),
                     sz=_TYPE.caption, color=_TXT)
                ry += Inches(h + 0.06)
            ry += Inches(0.08)
        if key_concepts and ry < Inches(5.9):
            _rect(sl, Inches(8.9), ry, Inches(4.1), Inches(0.25), _GRN)
            _txt(sl, 'LABELS TO KNOW', Inches(8.9), ry, Inches(4.1), Inches(0.25),
                 sz=_TYPE.chip, bold=True, color=_WHT, align=PP_ALIGN.CENTER)
            ry += Inches(0.32)
            for kc in key_concepts[:4]:
                if ry > Inches(6.4):
                    break
                lines = _wrapped_lines(f'• {kc}', 4.0, _TYPE.caption)
                h = 0.06 + lines * _line_h(_TYPE.caption)
                _txt(sl, f'• {kc}', Inches(8.95), ry, Inches(4.0), Inches(h),
                     sz=_TYPE.caption, bold=True, color=_ACC)
                ry += Inches(h + 0.04)
        if summary_txt:
            _txt(sl, _fit_line(f'↳ {summary_txt}', 8.2, _TYPE.footnote),
                 Inches(0.45), Inches(6.85),
                 Inches(8.2), Inches(0.4), sz=_TYPE.footnote, italic=True, color=_GRY)
        _add_slide_footer(sl, ks, content, is_dean=is_dean)

    # -- Renderer: generic fallback ---------------------------------------------------
    def _render_generic(sl, ks, content, *, is_dean: bool) -> None:
        stype = (content.get('slide_type') or '').upper()
        hcol  = {'SUMMARY': _GRN, 'QUIZ': _ORG, 'ACTIVITY': _TEAL}.get(stype, _NAV)
        _header(sl, ks.title or f'Slide {ks.slide_number}',
                subtitle=f'{course.code}  \xb7  Unit {kit.unit_number}  \xb7  Slide {ks.slide_number}',
                color=hcol)
        if stype:
            _rect(sl, Inches(11.15), Inches(0.08), Inches(1.95), Inches(0.27), _ACC)
            _txt(sl, stype, Inches(11.15), Inches(0.1), Inches(1.95), Inches(0.27),
                 sz=8, bold=True, color=_WHT, align=PP_ALIGN.CENTER)
        bullets        = [str(b) for b in (content.get('bullets') or []) if b]
        key_concepts   = [str(k) for k in (content.get('key_concepts') or []) if k]
        definitions    = [str(d) for d in (content.get('definitions') or []) if d]
        examples       = [str(e) for e in (content.get('examples') or []) if e]
        (content.get('code_snippet') or '').strip()
        activity       = (content.get('classroom_activity') or '').strip()
        summary_txt    = (content.get('student_summary') or '').strip()
        (content.get('teaching_notes') or '').strip()
        if not bullets:
            for _fk in ('points', 'body_points', 'slide_points', 'learning_points',
                        'content_points', 'teaching_points', 'body', 'slide_body'):
                _fv = content.get(_fk)
                if isinstance(_fv, list) and _fv:
                    bullets = [str(x) for x in _fv[:7] if x]; break
                elif isinstance(_fv, str) and _fv.strip():
                    bullets = [_fv.strip()]; break
        if not bullets and ks.speaker_notes:
            bullets = [ln.strip() for ln in ks.speaker_notes.split('\n') if ln.strip()][:5]
        if not bullets:
            if key_concepts:   bullets = ['Key concept: ' + kc for kc in key_concepts[:4]]
            elif definitions:  bullets = [d[:140] for d in definitions[:3]]
            elif examples:     bullets = ['Example: ' + ex for ex in examples[:3]]
            elif activity:     bullets = [activity[:140]]
            elif summary_txt:  bullets = [summary_txt[:140]]
            else:
                bullets = [f'Unit {kit.unit_number} - Slide {ks.slide_number}: '
                           'open Course Kit to add content before presenting']
        by = Inches(1.15)
        for bullet in bullets[:7]:
            if by > Inches(5.9):
                break
            bullet_str = '  ' + str(bullet)
            lines  = _wrapped_lines(bullet_str, 8.25, _TYPE.body)
            row_h  = max(0.5, 0.10 + lines * _line_h(_TYPE.body))
            _txt(sl, bullet_str, Inches(0.45), by,
                 Inches(8.25), Inches(row_h), sz=14, color=_TXT)
            by += Inches(row_h + 0.03)
        if activity and len(bullets) < 5:
            _rect(sl, Inches(0.45), by + Inches(0.1), Inches(8.25), Inches(0.26), _TEAL)
            _txt(sl, 'Activity: ' + activity, Inches(0.55), by + Inches(0.1),
                 Inches(8.1), Inches(0.26), sz=9, bold=True, color=_WHT)
        ry = Inches(1.15)
        if key_concepts:
            _rect(sl, Inches(8.9), ry, Inches(4.1), Inches(0.25), _ACC)
            _txt(sl, 'KEY CONCEPTS', Inches(8.9), ry,
                 Inches(4.1), Inches(0.25), sz=8, bold=True, color=_WHT, align=PP_ALIGN.CENTER)
            ry += Inches(0.3)
            for kc in key_concepts[:5]:
                _txt(sl, f'\u2022 {kc}', Inches(8.95), ry,
                     Inches(4.0), Inches(0.4), sz=11, bold=True, color=_ACC)
                ry += Inches(0.43)
            ry += Inches(0.1)
        if definitions:
            _rect(sl, Inches(8.9), ry, Inches(4.1), Inches(0.25), _GRN)
            _txt(sl, 'DEFINITIONS', Inches(8.9), ry,
                 Inches(4.1), Inches(0.25), sz=8, bold=True, color=_WHT, align=PP_ALIGN.CENTER)
            ry += Inches(0.3)
            for d in definitions[:3]:
                _txt(sl, f'\u27a4  {d}', Inches(8.95), ry,
                     Inches(4.0), Inches(0.52), sz=10, color=_TXT)
                ry += Inches(0.55)
        if examples:
            ey = Inches(6.15)
            _rect(sl, 0, ey, W, Inches(0.27), _TEAL)
            _txt(sl, 'EXAMPLES', Inches(0.4), ey, Inches(1.5), Inches(0.27),
                 sz=8, bold=True, color=_WHT)
            _txt(sl, '    \xb7    '.join(str(e) for e in examples[:3]),
                 Inches(2.1), ey, Inches(11), Inches(0.27), sz=9, italic=True, color=_WHT)
        if summary_txt:
            _txt(sl, _fit_line(f'\u21b3 {summary_txt}', 8.2, _TYPE.caption, max_lines=2),
                 Inches(0.4), Inches(6.5 if examples else 6.3),
                 Inches(8.2), Inches(0.65), sz=_TYPE.caption, italic=True, color=_GRY)
        _add_slide_footer(sl, ks, content, is_dean=is_dean)

    # -- Slide type registry ----------------------------------------------------------
    _SLIDE_REGISTRY = {
        'OBJECTIVES':      _render_objectives,
        'WORKED_EXAMPLE':  _render_worked_example,
        'CODE':            _render_code,
        'COMMON_MISTAKES': _render_common_mistakes,
        'ACTIVITY':        _render_activity,
        'QUIZ':            _render_quiz,
        'SUMMARY':         _render_summary_slide,
        'DIAGRAM':         _render_diagram,
    }

    # Slide types whose own layout has no room for a diagram: a quiz needs its
    # four options, a worked example its steps. Everything else defers to the
    # diagram when there is one to draw.
    _DIAGRAM_INCOMPATIBLE = {'QUIZ', 'CODE', 'WORKED_EXAMPLE', 'COMMON_MISTAKES', 'OBJECTIVES'}

    def _pick_renderer(content: dict, stype: str):
        """Dispatch on the diagram's PRESENCE, not only on the slide's label.

        The AI decides a slide has a diagram by emitting one; whether it also
        remembered to set slide_type='DIAGRAM' is a labelling detail, and a
        drawable diagram silently rendered as a bullet list is the exact bug
        this phase exists to fix.
        """
        if stype not in _DIAGRAM_INCOMPATIBLE and _parse_diagram(content) is not None:
            return _render_diagram
        return _SLIDE_REGISTRY.get(stype, _render_generic)

    # -- 3. CONTENT SLIDES ---------------------------------------------------------------
    for ks in slides_sorted:
        sl = prs.slides.add_slide(BLANK)
        _clear_placeholders(sl)
        _bg(sl, _WHT)
        content  = ks.content or {}
        stype    = (content.get('slide_type') or '').upper()
        renderer = _pick_renderer(content, stype)
        renderer(sl, ks, content, is_dean=is_dean)

    # â”€â”€ 4. TEACHING PLAN TABLE SLIDE â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€
    if tp:
        sl = prs.slides.add_slide(BLANK)
        _clear_placeholders(sl)
        _bg(sl, _WHT)
        _header(sl, 'Weekly Teaching Plan',
                subtitle=f'{course.code} â€” Unit {kit.unit_number}')
        n = min(len(tp), 10)
        tbl = sl.shapes.add_table(
            n + 1, 4, Inches(0.4), Inches(1.1), Inches(12.5), Inches(5.9)
        ).table
        _tbl_hdr(tbl, ('Wk', 'Topic', 'Objectives', 'CO Refs'),
                 (Inches(0.85), Inches(3.4), Inches(6.35), Inches(1.9)), _NAV)
        for ri, wk in enumerate(tp[:n], 1):
            _tbl_row(tbl, ri, [
                str(wk.get('week', ri)),
                str(wk.get('topic', ''))[:60],
                '  Â·  '.join(str(o) for o in (wk.get('objectives') or [])[:3])[:120],
                ', '.join(str(c2) for c2 in (wk.get('co_references') or [])),
            ], _LGRY)

    # â”€â”€ 5. SUMMARY SLIDE â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€
    sl = prs.slides.add_slide(BLANK)
    _clear_placeholders(sl)
    _bg(sl, _WHT)
    _header(sl, 'Unit Summary â€” Key Takeaways',
            subtitle=f'{course.code} â€” Unit {kit.unit_number}', color=_GRN)

    summaries = [(ks.content or {}).get('student_summary') or ''
                 for ks in slides_sorted]
    summaries = [s.strip() for s in summaries if s.strip()]

    items = summaries[:9] if summaries else \
            [f"{getattr(co,'code','')}: {getattr(co,'description','')}"
             for co in effective_cos[:9]]
    if items:
        _txt(sl, 'What students should know after this unit:',
             Inches(0.5), Inches(1.1), Inches(12.3), Inches(0.38),
             sz=12, bold=True, color=_ACC)
        for i, s in enumerate(items):
            _txt(sl, f'âœ“  {s}', Inches(0.6), Inches(1.58 + i * 0.58),
                 Inches(12.3), Inches(0.55), sz=13, color=_TXT)

    # â”€â”€ 6. RESOURCES SLIDE â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€
    resources = kit.resources or []
    if resources:
        sl = prs.slides.add_slide(BLANK)
        _clear_placeholders(sl)
        _bg(sl, _WHT)
        _header(sl, 'Teaching Resources & References',
                subtitle='Supplemental materials for faculty', color=_TEAL)
        ry2 = Inches(1.12)
        for res in resources[:10]:
            rtype = str(res.get('resource_type', '')).upper()[:12]
            title = str(res.get('title', ''))
            url   = str(res.get('url') or '')
            desc  = str(res.get('description') or '')
            _rect(sl, Inches(0.4), ry2, Inches(1.4), Inches(0.25), _ACC)
            _txt(sl, rtype, Inches(0.4), ry2, Inches(1.4), Inches(0.25),
                 sz=8, bold=True, color=_WHT, align=PP_ALIGN.CENTER)
            _txt(sl, title, Inches(1.9), ry2, Inches(11), Inches(0.28),
                 sz=12, bold=True, color=_TXT)
            detail = url or desc
            if detail:
                _txt(sl, detail[:110], Inches(1.9), ry2 + Inches(0.3),
                     Inches(11), Inches(0.25), sz=9, italic=True, color=_GRY)
                ry2 += Inches(0.68)
            else:
                ry2 += Inches(0.42)

    prs.save(buf)


# ---------------------------------------------------------------------------
# PPTX slide-body helper  (kept for PDF generation â€” not used in PPTX path)
# ---------------------------------------------------------------------------

def _fill_slide_body(tf, content: dict) -> None:
    """Render all H-34 SlideContent fields into a PPTX text-frame.

    teaching_notes is excluded here â€” it belongs in the notes pane only.
    Backward-compatible: missing or None fields are silently skipped.
    """
    from pptx.dml.color import RGBColor
    from pptx.util import Pt

    _ACCENT = RGBColor(0x44, 0x72, 0xC4)
    _GREEN  = RGBColor(0x21, 0x7A, 0x3C)
    _TEAL   = RGBColor(0x00, 0x7B, 0x8A)

    bullets            = content.get("bullets") or []
    key_concepts       = content.get("key_concepts") or []
    definitions        = content.get("definitions") or []
    examples           = content.get("examples") or []
    code_snippet       = (content.get("code_snippet") or "").strip()
    diagram_prompt     = (content.get("diagram_prompt") or "").strip()
    classroom_activity = (content.get("classroom_activity") or "").strip()
    student_summary    = (content.get("student_summary") or "").strip()

    tf.clear()
    _first = [True]

    def _para(text: str, *, level: int = 0, size: int = 18,
              bold: bool = False, italic: bool = False, color=None):
        if _first[0]:
            p = tf.paragraphs[0]
            _first[0] = False
        else:
            p = tf.add_paragraph()
        p.text  = text
        p.level = level
        for run in p.runs:
            run.font.size   = Pt(size)
            run.font.bold   = bold
            run.font.italic = italic
            if color is not None:
                run.font.color.rgb = color

    def _section(label: str, color=None):
        _para("", level=0, size=6)
        _para(label, level=0, size=13, bold=True, color=color or _ACCENT)

    for bullet in bullets:
        _para(bullet, level=0, size=18)

    if key_concepts:
        _section("Key Concepts")
        _para("  |  ".join(key_concepts), level=1, size=13, bold=True)

    if definitions:
        _section("Definitions", color=_GREEN)
        for defn in definitions:
            _para(f"â€¢ {defn}", level=1, size=14)

    if examples:
        _section("Examples", color=_TEAL)
        for ex in examples:
            _para(f"â€¢ {ex}", level=1, size=14)

    if code_snippet:
        _section("Code")
        for line in code_snippet.split("\n")[:20]:  # cap at 20 lines to prevent overflow
            _para(line or " ", level=1, size=11, italic=True)

    if diagram_prompt:
        _section("Diagram Prompt")
        _para(diagram_prompt, level=1, size=13)

    if classroom_activity:
        _section("Classroom Activity", color=_TEAL)
        _para(classroom_activity, level=1, size=13)

    if student_summary:
        _section("Student Summary")
        _para(student_summary, level=1, size=13)

# ---------------------------------------------------------------------------
# PDF generation (reportlab â€” slide handout + kit content)
# ---------------------------------------------------------------------------

def _generate_pdf(buf, kit, course, *, is_dean: bool) -> None:
    from reportlab.lib import colors
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import cm
    from reportlab.platypus import (
        HRFlowable,
        PageBreak,
        Paragraph,
        SimpleDocTemplate,
        Spacer,
        Table,
        TableStyle,
    )

    _BLUE   = colors.HexColor("#2E4057")
    _ACCENT = colors.HexColor("#4472C4")
    _ALT    = colors.HexColor("#EBF0FA")

    BASE_GRID = TableStyle([
        ("FONTNAME",      (0, 0), (-1, 0),  "Helvetica-Bold"),
        ("FONTNAME",      (0, 1), (-1, -1), "Helvetica"),
        ("FONTSIZE",      (0, 0), (-1, -1), 8),
        ("BACKGROUND",    (0, 0), (-1, 0),  _ACCENT),
        ("TEXTCOLOR",     (0, 0), (-1, 0),  colors.white),
        ("ROWBACKGROUNDS",(0, 1), (-1, -1), [colors.white, _ALT]),
        ("BOX",           (0, 0), (-1, -1), 0.5, colors.grey),
        ("INNERGRID",     (0, 0), (-1, -1), 0.25, colors.grey),
        ("VALIGN",        (0, 0), (-1, -1), "MIDDLE"),
        ("LEFTPADDING",   (0, 0), (-1, -1), 4),
        ("RIGHTPADDING",  (0, 0), (-1, -1), 4),
        ("TOPPADDING",    (0, 0), (-1, -1), 3),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 3),
    ])

    styles = getSampleStyleSheet()
    small  = ParagraphStyle("small", parent=styles["Normal"], fontSize=8, leading=10)
    tiny   = ParagraphStyle("tiny",  parent=styles["Normal"], fontSize=7, leading=9)
    h1     = ParagraphStyle("h1",    parent=styles["Heading1"], textColor=_BLUE)
    h2     = ParagraphStyle("h2",    parent=styles["Heading2"], textColor=_ACCENT)

    doc = SimpleDocTemplate(
        buf, pagesize=A4,
        topMargin=2*cm, bottomMargin=2*cm,
        leftMargin=2*cm, rightMargin=2*cm,
    )
    story = []

    # â”€â”€ Cover â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€
    story.append(Paragraph(f"{course.code}: {course.title}", styles["Title"]))
    story.append(Spacer(1, 0.2*cm))
    story.append(Paragraph(
        f"Unit {kit.unit_number} â€” Teaching Kit", styles["Heading2"]
    ))
    story.append(Spacer(1, 0.3*cm))
    story.append(HRFlowable(width="100%", thickness=2, color=_BLUE))
    story.append(Spacer(1, 0.4*cm))

    meta_rows = [
        ["Course Code",  course.code],
        ["Course Title", course.title],
        ["Unit",         str(kit.unit_number)],
        ["Kit Version",  str(kit.version)],
        ["Status",       kit.status.value],
        ["Complexity",   kit.complexity_level.value],
    ]
    if kit.published_at:
        meta_rows.append(["Published At", kit.published_at.strftime("%Y-%m-%d")])
    if kit.ai_model:
        meta_rows.append(["AI Model", kit.ai_model])

    meta_t = Table(meta_rows, colWidths=[4.5*cm, 10*cm])
    meta_t.setStyle(TableStyle([
        ("BACKGROUND",    (0, 0), (0, -1), colors.HexColor("#F2F2F2")),
        ("FONTNAME",      (0, 0), (-1, -1), "Helvetica"),
        ("FONTSIZE",      (0, 0), (-1, -1), 9),
        ("ROWBACKGROUNDS",(0, 0), (-1, -1), [colors.whitesmoke, colors.white]),
        ("BOX",           (0, 0), (-1, -1), 0.5, colors.grey),
        ("INNERGRID",     (0, 0), (-1, -1), 0.25, colors.grey),
        ("VALIGN",        (0, 0), (-1, -1), "MIDDLE"),
        ("LEFTPADDING",   (0, 0), (-1, -1), 6),
        ("TOPPADDING",    (0, 0), (-1, -1), 4),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 4),
    ]))
    story.append(meta_t)
    story.append(Spacer(1, 1*cm))

    # â”€â”€ Slides â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€
    slides_sorted = sorted(kit.slides or [], key=lambda s: s.slide_number)
    if slides_sorted:
        story.append(Paragraph("Slides", h1))
        story.append(Spacer(1, 0.3*cm))

        for kit_slide in slides_sorted:
            story.append(Paragraph(
                f"Slide {kit_slide.slide_number}: {kit_slide.title or '(untitled)'}",
                h2,
            ))
            content      = kit_slide.content or {}
            bullets      = content.get("bullets") or []
            key_concepts = content.get("key_concepts") or []

            for bullet in bullets:
                story.append(Paragraph(f"â€¢ {bullet}", small))
            if key_concepts:
                story.append(Paragraph(
                    "Key concepts: " + "  |  ".join(key_concepts), tiny
                ))

            footer_parts = []
            if kit_slide.bloom_level:
                footer_parts.append(f"Bloom: {kit_slide.bloom_level.value}")
            if kit_slide.co_reference:
                footer_parts.append(f"CO: {kit_slide.co_reference}")
            if footer_parts:
                story.append(Paragraph(" | ".join(footer_parts), tiny))
            if not is_dean and kit_slide.speaker_notes:
                story.append(Paragraph(
                    f"[Notes] {kit_slide.speaker_notes}", tiny
                ))
            story.append(Spacer(1, 0.4*cm))

        story.append(PageBreak())

    # â”€â”€ Weekly teaching plan â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€
    teaching_plan = kit.teaching_plan or []
    if teaching_plan:
        story.append(Paragraph("Weekly Teaching Plan", h1))
        story.append(Spacer(1, 0.3*cm))
        header = [["Week", "Topic", "Hours", "CO References"]]
        rows = []
        for week in teaching_plan:
            co_refs = ", ".join(week.get("co_references") or []) or "â€”"
            rows.append([
                str(week.get("week", "?")),
                Paragraph(week.get("topic", ""), small),
                str(week.get("hours", "?")),
                Paragraph(co_refs, small),
            ])
        t = Table(header + rows, colWidths=[1.5*cm, 8*cm, 2*cm, 4*cm])
        t.setStyle(BASE_GRID)
        story.append(t)
        story.append(Spacer(1, 1*cm))

    # â”€â”€ Assignments â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€
    assignments_sorted = sorted(
        kit.assignments or [], key=lambda a: a.assignment_number
    )
    if assignments_sorted:
        story.append(Paragraph("Assignments", h1))
        story.append(Spacer(1, 0.3*cm))
        for asn in assignments_sorted:
            story.append(Paragraph(
                f"Assignment {asn.assignment_number}: {asn.title}", h2
            ))
            story.append(Paragraph(
                f"Type: {asn.assignment_type.value}  |  "
                f"Complexity: {asn.complexity_level.value}",
                tiny,
            ))
            story.append(Paragraph(asn.question_text, small))
            if not is_dean and asn.model_answer:
                story.append(Paragraph(f"Model Answer: {asn.model_answer}", small))
            rubric = asn.rubric or []
            if rubric and not is_dean:
                story.append(Paragraph("Rubric:", small))
                rub_h = [["Criterion", "Description", "Marks"]]
                rub_rows = [
                    [
                        Paragraph(r.get("criterion", ""), small),
                        Paragraph(r.get("description", ""), small),
                        str(r.get("max_marks", "")),
                    ]
                    for r in rubric if isinstance(r, dict)
                ]
                rub_t = Table(rub_h + rub_rows, colWidths=[3.5*cm, 9.5*cm, 2.5*cm])
                rub_t.setStyle(BASE_GRID)
                story.append(rub_t)
            story.append(Spacer(1, 0.5*cm))

    # â”€â”€ Teaching Resources â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€
    resources = kit.resources or []
    if resources:
        story.append(Paragraph("Teaching Resources", h1))
        story.append(Spacer(1, 0.3*cm))
        for res in resources:
            rtype     = res.get("resource_type", "")
            title_str = res.get("title", "")
            url       = res.get("url")
            desc      = res.get("description")
            story.append(Paragraph(f"â€¢ [{rtype}] {title_str}", small))
            if desc:
                story.append(Paragraph(f"  {desc}", tiny))
            if url:
                story.append(Paragraph(f"  {url}", tiny))

    doc.build(story)


# ---------------------------------------------------------------------------
# Student handout PDF (sanitized â€” no speaker_notes, answer_key, model_answer, rubric)
# ---------------------------------------------------------------------------

def _generate_handout_pdf(buf, kit, course) -> None:
    """
    Produces a student-facing A4 PDF with a diagonal watermark.
    All faculty-sensitive fields are excluded:
      - KitSlide.speaker_notes
      - KitAssignment.model_answer and rubric
    Available to ADMIN, FACULTY, and DEAN roles.
    """
    from reportlab.lib import colors
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import cm
    from reportlab.platypus import (
        HRFlowable, PageBreak, Paragraph, SimpleDocTemplate, Spacer, Table, TableStyle,
    )

    _BLUE   = colors.HexColor("#2E4057")
    _ACCENT = colors.HexColor("#4472C4")
    _ALT    = colors.HexColor("#EBF0FA")
    _WMARK  = colors.Color(0.80, 0.80, 0.80, alpha=0.30)

    watermark_text = f"STUDENT HANDOUT â€” {course.code}"

    def _draw_watermark(canv, doc):
        canv.saveState()
        canv.setFont("Helvetica", 36)
        canv.setFillColor(_WMARK)
        canv.translate(A4[0] / 2, A4[1] / 2)
        canv.rotate(45)
        canv.drawCentredString(0, 0, watermark_text)
        canv.restoreState()

    TableStyle([
        ("FONTNAME",      (0, 0), (-1, 0),  "Helvetica-Bold"),
        ("FONTNAME",      (0, 1), (-1, -1), "Helvetica"),
        ("FONTSIZE",      (0, 0), (-1, -1), 8),
        ("BACKGROUND",    (0, 0), (-1, 0),  _ACCENT),
        ("TEXTCOLOR",     (0, 0), (-1, 0),  colors.white),
        ("ROWBACKGROUNDS",(0, 1), (-1, -1), [colors.white, _ALT]),
        ("BOX",           (0, 0), (-1, -1), 0.5, colors.grey),
        ("INNERGRID",     (0, 0), (-1, -1), 0.25, colors.grey),
        ("VALIGN",        (0, 0), (-1, -1), "MIDDLE"),
        ("LEFTPADDING",   (0, 0), (-1, -1), 4),
        ("RIGHTPADDING",  (0, 0), (-1, -1), 4),
        ("TOPPADDING",    (0, 0), (-1, -1), 3),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 3),
    ])

    styles = getSampleStyleSheet()
    ho_small = ParagraphStyle("ho_small", parent=styles["Normal"], fontSize=8,  leading=10)
    ho_tiny  = ParagraphStyle("ho_tiny",  parent=styles["Normal"], fontSize=7,  leading=9)
    ho_h1    = ParagraphStyle("ho_h1",    parent=styles["Heading1"], textColor=_BLUE)
    ho_h2    = ParagraphStyle("ho_h2",    parent=styles["Heading2"], textColor=_ACCENT)

    doc = SimpleDocTemplate(
        buf, pagesize=A4,
        topMargin=2*cm, bottomMargin=2*cm,
        leftMargin=2*cm, rightMargin=2*cm,
    )
    story = []

    # â”€â”€ Cover â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€
    story.append(Paragraph(f"{course.code}: {course.title}", styles["Title"]))
    story.append(Spacer(1, 0.2*cm))
    story.append(Paragraph(f"Unit {kit.unit_number} â€” Student Handout", styles["Heading2"]))
    story.append(Spacer(1, 0.3*cm))
    story.append(HRFlowable(width="100%", thickness=2, color=_BLUE))
    story.append(Spacer(1, 0.4*cm))

    meta_rows = [
        ["Course Code",  course.code],
        ["Course Title", course.title],
        ["Unit",         str(kit.unit_number)],
        ["Complexity",   kit.complexity_level.value],
    ]
    if kit.published_at:
        meta_rows.append(["Published At", kit.published_at.strftime("%Y-%m-%d")])

    meta_t = Table(meta_rows, colWidths=[4.5*cm, 10*cm])
    meta_t.setStyle(TableStyle([
        ("BACKGROUND",    (0, 0), (0, -1), colors.HexColor("#F2F2F2")),
        ("FONTNAME",      (0, 0), (-1, -1), "Helvetica"),
        ("FONTSIZE",      (0, 0), (-1, -1), 9),
        ("ROWBACKGROUNDS",(0, 0), (-1, -1), [colors.whitesmoke, colors.white]),
        ("BOX",           (0, 0), (-1, -1), 0.5, colors.grey),
        ("INNERGRID",     (0, 0), (-1, -1), 0.25, colors.grey),
        ("VALIGN",        (0, 0), (-1, -1), "MIDDLE"),
        ("LEFTPADDING",   (0, 0), (-1, -1), 6),
        ("TOPPADDING",    (0, 0), (-1, -1), 4),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 4),
    ]))
    story.append(meta_t)
    story.append(Spacer(1, 1*cm))

    # â”€â”€ Slides (content only â€” no speaker_notes) â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€
    slides_sorted = sorted(kit.slides or [], key=lambda s: s.slide_number)
    if slides_sorted:
        story.append(Paragraph("Slides", ho_h1))
        story.append(Spacer(1, 0.3*cm))
        for kit_slide in slides_sorted:
            story.append(Paragraph(
                f"Slide {kit_slide.slide_number}: {kit_slide.title or '(untitled)'}",
                ho_h2,
            ))
            content      = kit_slide.content or {}
            bullets      = content.get("bullets") or []
            key_concepts = content.get("key_concepts") or []
            for bullet in bullets:
                story.append(Paragraph(f"â€¢ {bullet}", ho_small))
            if key_concepts:
                story.append(Paragraph(
                    "Key concepts: " + "  |  ".join(key_concepts), ho_tiny,
                ))
            footer_parts = []
            if kit_slide.bloom_level:
                footer_parts.append(f"Bloom: {kit_slide.bloom_level.value}")
            if kit_slide.co_reference:
                footer_parts.append(f"CO: {kit_slide.co_reference}")
            if footer_parts:
                story.append(Paragraph(" | ".join(footer_parts), ho_tiny))
            story.append(Spacer(1, 0.4*cm))
        story.append(PageBreak())

    # â”€â”€ Assignments (question text only â€” no model_answer, no rubric) â”€â”€â”€â”€â”€
    assignments_sorted = sorted(
        kit.assignments or [], key=lambda a: a.assignment_number
    )
    if assignments_sorted:
        story.append(Paragraph("Assignments", ho_h1))
        story.append(Spacer(1, 0.3*cm))
        for asn in assignments_sorted:
            story.append(Paragraph(
                f"Assignment {asn.assignment_number}: {asn.title}", ho_h2,
            ))
            story.append(Paragraph(
                f"Type: {asn.assignment_type.value}  |  "
                f"Complexity: {asn.complexity_level.value}",
                ho_tiny,
            ))
            story.append(Paragraph(asn.question_text, ho_small))
            story.append(Spacer(1, 0.5*cm))

    # â”€â”€ Teaching resources â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€â”€
    resources = kit.resources or []
    if resources:
        story.append(Paragraph("Teaching Resources", ho_h1))
        story.append(Spacer(1, 0.3*cm))
        for res in resources:
            rtype     = res.get("resource_type", "")
            title_str = res.get("title", "")
            url       = res.get("url")
            desc      = res.get("description")
            story.append(Paragraph(f"â€¢ [{rtype}] {title_str}", ho_small))
            if desc:
                story.append(Paragraph(f"  {desc}", ho_tiny))
            if url:
                story.append(Paragraph(f"  {url}", ho_tiny))

    doc.build(story, onFirstPage=_draw_watermark, onLaterPages=_draw_watermark)
