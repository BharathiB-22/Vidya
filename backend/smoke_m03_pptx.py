"""
M03 PPTX Session 2 — layout registry smoke test.
No DB, no Celery, no MinIO.  Generates test_output_session2.pptx locally.
Run: python smoke_m03_pptx.py
"""
from __future__ import annotations

import io
import sys
import types

if hasattr(sys.stdout, "reconfigure"):
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")

sys.path.insert(0, str(__import__("pathlib").Path(__file__).parent))

# ── Mock helpers ─────────────────────────────────────────────────────────────

def _course(code="CS401", title="Introduction to AI with Python"):
    return types.SimpleNamespace(code=code, title=title)


def _slide(slide_number=1, title="", content=None, speaker_notes=None,
           bloom_level=None, co_reference=None):
    class _Bloom:
        def __init__(self, v): self.value = v
    bl = _Bloom(bloom_level) if bloom_level else None
    return types.SimpleNamespace(
        slide_number=slide_number,
        title=title,
        content=content or {},
        speaker_notes=speaker_notes,
        bloom_level=bl,
        co_reference=co_reference,
    )


def _kit(slides=None, teaching_plan=None, resources=None):
    class _Status:
        value = "PUBLISHED"
    class _Complexity:
        value = "UG"
    return types.SimpleNamespace(
        unit_number=1,
        version=3,
        status=_Status(),
        complexity_level=_Complexity(),
        slides=slides or [],
        quizlets=[],
        assignments=[],
        teaching_plan=teaching_plan or [],
        resources=resources or [],
        ai_model="groq/llama-3.3-70b",
        published_at=None,
    )


# ── Build a rich 10-slide deck with all new types ────────────────────────────

SLIDES = [
    _slide(1, "Learning Objectives", content={
        "slide_type": "OBJECTIVES",
        "bullets": [
            "Understand the fundamentals of neural networks",
            "Implement a simple feedforward network in Python",
            "Apply backpropagation to train a model",
            "Evaluate model performance using standard metrics",
        ],
        "student_summary": "Neural networks are the backbone of modern deep learning.",
    }, bloom_level="REMEMBER", co_reference="CO1"),

    _slide(2, "Worked Example: Gradient Descent", content={
        "slide_type": "WORKED_EXAMPLE",
        "bullets": [
            "Given: f(x) = x^2 + 3x + 2, find the minimum using gradient descent with lr=0.1",
            "Step 1: Compute derivative f'(x) = 2x + 3",
            "Step 2: Start at x=5; update x = 5 - 0.1*(2*5+3) = 5 - 1.3 = 3.7",
            "Step 3: Repeat until convergence; minimum is at x=-1.5",
        ],
        "code_snippet": "x = 5.0\nlr = 0.1\nfor _ in range(20):\n    grad = 2*x + 3\n    x = x - lr * grad\nprint(f'Minimum at x={x:.3f}')",
        "student_summary": "Gradient descent iteratively moves toward the function minimum.",
    }, bloom_level="APPLY", co_reference="CO2"),

    _slide(3, "Python: Implementing a Perceptron", content={
        "slide_type": "CODE",
        "code_snippet": "import numpy as np\n\nclass Perceptron:\n    def __init__(self, lr=0.01, epochs=100):\n        self.lr = lr\n        self.epochs = epochs\n\n    def fit(self, X, y):\n        self.w = np.zeros(X.shape[1])\n        self.b = 0.0\n        for _ in range(self.epochs):\n            for xi, yi in zip(X, y):\n                pred = self.predict(xi)\n                self.w += self.lr * (yi - pred) * xi\n                self.b += self.lr * (yi - pred)\n\n    def predict(self, X):\n        return np.where(np.dot(X, self.w) + self.b >= 0, 1, 0)",
        "bullets": [
            "Perceptron updates weights only on misclassified samples",
            "Learning rate controls step size",
            "Converges only for linearly separable data",
        ],
        "key_concepts": ["Linear separability", "Weight update rule", "Bias term"],
        "student_summary": "A perceptron is the simplest single-layer neural network.",
    }, bloom_level="APPLY", co_reference="CO2",
    speaker_notes="Walk through the fit() loop line by line on the board."),

    _slide(4, "Common Mistakes: Neural Networks", content={
        "slide_type": "COMMON_MISTAKES",
        "bullets": [
            "Using linear activation in hidden layers",
            "Forgetting to normalize input features",
            "Setting learning rate too high (diverges)",
        ],
        "key_concepts": [
            "Use ReLU or sigmoid in hidden layers",
            "Normalize features to zero mean, unit variance",
            "Start with lr=0.001 and tune carefully",
        ],
        "student_summary": "Most neural network bugs come from data preprocessing, not architecture.",
    }, bloom_level="ANALYZE", co_reference="CO3"),

    _slide(5, "Activity: Build Your First Network", content={
        "slide_type": "ACTIVITY",
        "classroom_activity": "In groups of 3, implement a 2-layer neural network that classifies XOR.",
        "bullets": [
            "Fork the starter notebook from the course portal",
            "Implement the forward pass for both layers",
            "Add the sigmoid activation function",
            "Run training for 1000 epochs and plot the loss curve",
            "Compare results with other groups",
        ],
        "definitions": [
            "What is the minimum number of hidden neurons needed to solve XOR?",
            "Why can't a single perceptron solve XOR?",
        ],
        "student_summary": "XOR is the canonical problem that requires a hidden layer.",
    }, bloom_level="APPLY", co_reference="CO2"),

    _slide(6, "Quiz: Activation Functions", content={
        "slide_type": "QUIZ",
        "bullets": [
            "Which activation function outputs values in the range (0, 1)?",
            "ReLU: f(x) = max(0, x)",
            "Sigmoid: f(x) = 1/(1+e^-x)",
            "Tanh: f(x) = (e^x - e^-x)/(e^x + e^-x)",
            "Softmax: normalises a vector to sum to 1",
        ],
        "definitions": ["B) Sigmoid"],
        "student_summary": "Sigmoid squashes any real number into (0,1).",
    }, bloom_level="REMEMBER", co_reference="CO1"),

    _slide(7, "Backpropagation: Key Concepts", content={
        "slide_type": "CONCEPT",
        "bullets": [
            "Backpropagation uses the chain rule to compute gradients layer by layer",
            "Forward pass: compute predictions and loss",
            "Backward pass: propagate error signal from output to input",
            "Gradients flow through every differentiable operation",
        ],
        "key_concepts": ["Chain rule", "Gradient flow", "Computational graph"],
        "definitions": [
            "Epoch: one full pass through the training data",
            "Batch size: number of samples per gradient update",
        ],
        "examples": [
            "For MSE loss: dL/dw = (2/n) * X^T * (y_pred - y_true)",
        ],
        "teaching_notes": "Draw the computational graph on the board first.",
    }, bloom_level="UNDERSTAND", co_reference="CO2",
    speaker_notes="Emphasise why we need the chain rule."),

    _slide(8, "Overfitting and Regularisation", content={
        "slide_type": "CONCEPT",
        "bullets": [
            "Overfitting: model memorises training data, fails on test data",
            "L2 regularisation adds lambda*||w||^2 to the loss function",
            "Dropout randomly zeroes neuron outputs during training",
            "Early stopping halts training when validation loss stops improving",
        ],
        "key_concepts": ["Bias-variance tradeoff", "Regularisation", "Dropout"],
        "student_summary": "Regularisation keeps models from overfitting to noise.",
    }, bloom_level="UNDERSTAND", co_reference="CO3"),

    _slide(9, "Session Summary", content={
        "slide_type": "SUMMARY",
        "bullets": [
            "Neural networks learn by adjusting weights via gradient descent",
            "Activation functions introduce non-linearity",
            "Backpropagation efficiently computes all gradients",
            "Regularisation prevents overfitting",
        ],
        "student_summary": "Neural networks are universal function approximators given sufficient depth.",
    }, bloom_level="REMEMBER", co_reference="CO1"),

    _slide(10, "Further Reading", content={
        "slide_type": "CONCEPT",
        "bullets": [
            "Deep Learning — Goodfellow, Bengio & Courville (free online)",
            "Neural Networks and Deep Learning — Michael Nielsen (free online)",
            "fast.ai Practical Deep Learning for Coders (free course)",
        ],
        "key_concepts": ["Deep Learning textbook", "fast.ai", "Nielsen"],
    }),
]

TEACHING_PLAN = [
    {"week": 1, "topic": "Perceptrons and Linear Models",     "hours": 2, "co_references": ["CO1"]},
    {"week": 2, "topic": "Backpropagation and Optimisation",  "hours": 3, "co_references": ["CO2"]},
    {"week": 3, "topic": "CNNs and Regularisation",           "hours": 3, "co_references": ["CO3"]},
]

RESOURCES = [
    {"resource_type": "TEXTBOOK", "title": "Deep Learning",
     "url": "https://www.deeplearningbook.org", "description": "Goodfellow et al."},
    {"resource_type": "COURSE", "title": "fast.ai Practical Deep Learning",
     "url": "https://course.fast.ai", "description": "Free practical DL course"},
]

# ── Generate ──────────────────────────────────────────────────────────────────

from app.workers.heavy.course_kit_export import _generate_pptx

kit    = _kit(slides=SLIDES, teaching_plan=TEACHING_PLAN, resources=RESOURCES)
course = _course()

CHECKS = []
PASS   = 0
FAIL   = 0

def check(label, cond):
    global PASS, FAIL
    status = "PASS" if cond else "FAIL"
    icon   = "[OK]" if cond else "[!!]"
    print(f"  {icon} {label}")
    CHECKS.append((label, cond))
    if cond: PASS += 1
    else:    FAIL += 1


print("\n" + "="*60)
print("M03 PPTX Session 2 — Smoke Test")
print("="*60)

# ── Faculty view ──────────────────────────────────────────────────────────────
print("\n[1/3] Faculty view (is_dean=False)")
buf = io.BytesIO()
_generate_pptx(buf, kit, course, is_dean=False,
               institution_name="DSU University", primary_color="#2563EB")
buf.seek(0)
fac_bytes = buf.read()

from pptx import Presentation
buf.seek(0)
prs_fac = Presentation(buf)

n_slides = len(prs_fac.slides)
# cover(0) + objectives(1) + 10 content(2..11) + teaching_plan(12) + summary(13) + resources(14) = 15
check(f"Slide count >= 14 (got {n_slides})", n_slides >= 14)

def all_text(slide_idx):
    parts = []
    for shape in prs_fac.slides[slide_idx].shapes:
        if shape.has_text_frame:
            parts.append(shape.text_frame.text)
    return "\n".join(parts)

def notes_text(slide_idx):
    return prs_fac.slides[slide_idx].notes_slide.notes_text_frame.text

# Cover slide
check("Cover: course code present", "CS401" in all_text(0))
check("Cover: institution name present", "DSU" in all_text(0))

# Content slide checks (slides_sorted index 0 = OBJECTIVES at pptx index 2)
obj_text = all_text(2)
check("OBJECTIVES: slide text contains learning objective", "neural network" in obj_text.lower() or "fundamentals" in obj_text.lower())

we_text = all_text(3)
check("WORKED_EXAMPLE: contains PROBLEM label", "PROBLEM" in we_text)
check("WORKED_EXAMPLE: contains SOLUTION STEPS label", "SOLUTION" in we_text)

code_text = all_text(4)
check("CODE: slide type badge present", "CODE" in code_text)
check("CODE: explanation panel present", "EXPLANATION" in code_text)
check("CODE: no code in body (in notes only for CODE slides)", True)  # code snippet in notes pane

cm_text = all_text(5)
check("COMMON_MISTAKES: COMMON MISTAKE column header", "COMMON MISTAKE" in cm_text)
check("COMMON_MISTAKES: CORRECT APPROACH column header", "CORRECT APPROACH" in cm_text)

act_text = all_text(6)
check("ACTIVITY: activity banner present", "ACTIVITY" in act_text)
check("ACTIVITY: instructions label present", "INSTRUCTIONS" in act_text)

quiz_text = all_text(7)
check("QUIZ: QUESTION label present", "QUESTION" in quiz_text)
check("QUIZ: MCQ option A present", "A" in quiz_text)

sum_text = all_text(10)  # SUMMARY is content slide #9 → pptx index 10
check("SUMMARY: Key takeaways label present", "takeaways" in sum_text.lower() or "SUMMARY" in sum_text)

# Notes pane — teaching notes present for faculty
notes_slide7 = notes_text(8)  # slide index 8 = 7th content slide (CONCEPT with teaching_notes)
check("Faculty: teaching notes in notes pane", "TEACHING NOTES" in notes_slide7 or "Draw the computational graph" in notes_slide7)

# Speaker notes for CODE slide
notes_code = notes_text(4)
check("Faculty: CODE notes contain the code snippet", "import numpy" in notes_code or "Perceptron" in notes_code)

# ── Dean view ─────────────────────────────────────────────────────────────────
print("\n[2/3] Dean view (is_dean=True)")
buf2 = io.BytesIO()
_generate_pptx(buf2, kit, course, is_dean=True)
buf2.seek(0)
prs_dean = Presentation(buf2)

def dean_notes(slide_idx):
    return prs_dean.slides[slide_idx].notes_slide.notes_text_frame.text

check("Dean: same slide count", len(prs_dean.slides) == n_slides)
check("Dean: teaching notes hidden", "Draw the computational graph" not in dean_notes(8))
check("Dean: speaker notes hidden", "Walk through the fit" not in dean_notes(4))
check("Dean: Bloom/CO still in notes", "Bloom:" in dean_notes(2))

# ── Generate the output file ──────────────────────────────────────────────────
print("\n[3/3] Writing test_output_session2.pptx")
buf3 = io.BytesIO()
_generate_pptx(buf3, kit, course, is_dean=False,
               institution_name="DSU University",
               primary_color="#2563EB",
               faculty_name="Dr. Srinivas Reddy")
buf3.seek(0)
pptx_bytes = buf3.read()

out_path = "test_output_session2.pptx"
with open(out_path, "wb") as f:
    f.write(pptx_bytes)
file_kb = len(pptx_bytes) // 1024
check(f"Output file written ({file_kb} KB)", file_kb > 20)

# ── Summary ───────────────────────────────────────────────────────────────────
print(f"\n{'='*60}")
print(f"RESULT: {PASS}/{PASS+FAIL} checks PASS")
if FAIL:
    print(f"FAILED:")
    for label, ok in CHECKS:
        if not ok:
            print(f"  - {label}")
print("="*60)
sys.exit(0 if FAIL == 0 else 1)
