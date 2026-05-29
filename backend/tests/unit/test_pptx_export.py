"""
Unit tests for course_kit_export._generate_pptx.

Verifies that:
  - Every content slide has at least one text-bearing body textbox.
  - bullets, key_concepts, definitions, examples are rendered when present.
  - The three-level fallback (alt field names → speaker_notes → synthesise)
    ensures a slide body is never blank.
  - No "Click to add text" placeholder XML survives in exported slides.

No network, no DB, no Celery required.
"""
import io
import os
import types

import pytest

# Minimal env setup so app.config loads without real credentials.
os.environ.setdefault("DATABASE_URL", "postgresql+asyncpg://x:x@localhost/x")
os.environ.setdefault("JWT_SECRET", "x")
os.environ.setdefault("GEMINI_API_KEY", "")
os.environ.setdefault("GROQ_API_KEY", "placeholder")
os.environ.setdefault("S3_ENDPOINT", "http://localhost:9000")
os.environ.setdefault("S3_ACCESS_KEY", "x")
os.environ.setdefault("S3_SECRET_KEY", "x")


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _make_slide(
    n: int,
    title: str,
    *,
    bullets=None,
    key_concepts=None,
    definitions=None,
    examples=None,
    slide_type: str = "CONCEPT",
    speaker_notes: str | None = None,
    bloom_level: str | None = None,
    co_reference: str | None = None,
    classroom_activity: str | None = None,
    student_summary: str | None = None,
    teaching_notes: str | None = None,
    code_snippet: str | None = None,
    alt_field: dict | None = None,   # inject arbitrary extra fields into content
):
    """Return a SimpleNamespace that mimics a KitSlide ORM object."""
    content: dict = {
        "slide_type":         slide_type,
        "bullets":            bullets or [],
        "key_concepts":       key_concepts or [],
        "definitions":        definitions or [],
        "examples":           examples or [],
        "classroom_activity": classroom_activity,
        "student_summary":    student_summary,
        "teaching_notes":     teaching_notes,
        "code_snippet":       code_snippet,
    }
    if alt_field:
        content.update(alt_field)

    bloom = None
    if bloom_level:
        bloom = types.SimpleNamespace(value=bloom_level)

    return types.SimpleNamespace(
        slide_number=n,
        title=title,
        content=content,
        speaker_notes=speaker_notes,
        bloom_level=bloom,
        co_reference=co_reference,
    )


def _make_kit(slides, *, unit_number: int = 1, version: int = 1,
              teaching_plan=None, resources=None):
    """Return a SimpleNamespace that mimics a CourseKit ORM object."""
    # Minimal enum-like for complexity_level
    complexity = types.SimpleNamespace(value="UG")
    return types.SimpleNamespace(
        unit_number=unit_number,
        version=version,
        slides=slides,
        quizlets=[],
        assignments=[],
        teaching_plan=teaching_plan or [],
        lesson_plans=[],
        resources=resources or [],
        complexity_level=complexity,
        published_at=None,
        ai_model="test-model",
    )


def _make_course(code: str = "CS101", title: str = "Introduction to Computer Science"):
    return types.SimpleNamespace(code=code, title=title)


def _run_pptx(kit, course) -> "pptx.Presentation":
    """Call _generate_pptx and return the parsed Presentation object."""
    from pptx import Presentation as _Prs  # noqa: F401 — import check
    from app.workers.heavy.course_kit_export import _generate_pptx

    buf = io.BytesIO()
    _generate_pptx(buf, kit, course, is_dean=False)
    buf.seek(0)
    from pptx import Presentation
    return Presentation(buf)


def _body_texts(slide) -> list[str]:
    """Return all non-empty text strings from a slide's text shapes (excluding rects)."""
    texts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            t = shape.text_frame.text.strip()
            if t:
                texts.append(t)
    return texts


def _slide_index_of_content_slides(prs) -> list[int]:
    """Return 0-based indices of content (non-cover) slides."""
    # slide 0 = cover, slide 1 = objectives — skip those two
    return list(range(2, len(prs.slides)))


# ---------------------------------------------------------------------------
# Tests
# ---------------------------------------------------------------------------

class TestPptxBullets:
    def test_bullets_rendered_in_body(self):
        """Content slide with bullets must have those texts in the exported PPTX."""
        slide_obj = _make_slide(
            1, "Python Variables",
            bullets=[
                "Variables store data in memory",
                "Python uses dynamic typing",
                "Common types: int, str, list, dict",
            ],
        )
        kit  = _make_kit([slide_obj])
        prs  = _run_pptx(kit, _make_course())
        # Content slides start at index 2
        content_slides = [prs.slides[i] for i in _slide_index_of_content_slides(prs)]
        assert len(content_slides) >= 1
        all_texts = _body_texts(content_slides[0])
        combined = " ".join(all_texts)
        assert "Variables store data in memory" in combined
        assert "Python uses dynamic typing" in combined

    def test_key_concepts_rendered_in_right_panel(self):
        """key_concepts must appear as text even when bullets are also present."""
        slide_obj = _make_slide(
            1, "OOP Concepts",
            bullets=["Classes define blueprints for objects"],
            key_concepts=["encapsulation", "inheritance", "polymorphism"],
        )
        kit = _make_kit([slide_obj])
        prs = _run_pptx(kit, _make_course())
        content_slides = [prs.slides[i] for i in _slide_index_of_content_slides(prs)]
        all_texts = _body_texts(content_slides[0])
        combined = " ".join(all_texts)
        assert "encapsulation" in combined.lower()

    def test_definitions_rendered(self):
        """definitions must appear in the exported slide."""
        slide_obj = _make_slide(
            1, "Key Definitions",
            bullets=["Study these definitions carefully"],
            definitions=[
                "Algorithm: A step-by-step procedure to solve a problem",
                "Recursion: A function that calls itself",
            ],
        )
        kit = _make_kit([slide_obj])
        prs = _run_pptx(kit, _make_course())
        content_slides = [prs.slides[i] for i in _slide_index_of_content_slides(prs)]
        all_texts = _body_texts(content_slides[0])
        combined = " ".join(all_texts)
        assert "Algorithm" in combined

    def test_examples_rendered(self):
        """examples must appear in the exported slide."""
        slide_obj = _make_slide(
            1, "Real-World Examples",
            bullets=["Applications of sorting algorithms"],
            examples=["QuickSort used in Python's sorted()", "MergeSort in Java's Arrays.sort()"],
        )
        kit = _make_kit([slide_obj])
        prs = _run_pptx(kit, _make_course())
        content_slides = [prs.slides[i] for i in _slide_index_of_content_slides(prs)]
        all_texts = _body_texts(content_slides[0])
        combined = " ".join(all_texts)
        assert "QuickSort" in combined


class TestPptxFallbacks:
    def test_empty_bullets_falls_back_to_key_concepts(self):
        """When bullets is empty but key_concepts has content, body is not blank."""
        slide_obj = _make_slide(
            1, "Core Concepts",
            bullets=[],
            key_concepts=["abstraction", "modularity", "reusability"],
        )
        kit = _make_kit([slide_obj])
        prs = _run_pptx(kit, _make_course())
        content_slides = [prs.slides[i] for i in _slide_index_of_content_slides(prs)]
        all_texts = _body_texts(content_slides[0])
        combined = " ".join(all_texts).lower()
        # Fallback generates "Key concept: abstraction" etc.
        assert "abstraction" in combined

    def test_empty_bullets_falls_back_to_definitions(self):
        """When bullets AND key_concepts empty, definitions are used."""
        slide_obj = _make_slide(
            1, "Definitions Slide",
            bullets=[],
            key_concepts=[],
            definitions=["Stack: LIFO data structure", "Queue: FIFO data structure"],
        )
        kit = _make_kit([slide_obj])
        prs = _run_pptx(kit, _make_course())
        content_slides = [prs.slides[i] for i in _slide_index_of_content_slides(prs)]
        all_texts = _body_texts(content_slides[0])
        combined = " ".join(all_texts)
        assert "Stack" in combined

    def test_empty_bullets_falls_back_to_examples(self):
        """When bullets/key_concepts/definitions empty, examples are used."""
        slide_obj = _make_slide(
            1, "Examples Only",
            bullets=[], key_concepts=[], definitions=[],
            examples=["Binary search on a sorted list"],
        )
        kit = _make_kit([slide_obj])
        prs = _run_pptx(kit, _make_course())
        content_slides = [prs.slides[i] for i in _slide_index_of_content_slides(prs)]
        all_texts = _body_texts(content_slides[0])
        combined = " ".join(all_texts)
        assert "Binary search" in combined

    def test_empty_all_falls_back_to_speaker_notes(self):
        """When all content fields empty, speaker_notes are split into bullets."""
        slide_obj = _make_slide(
            1, "Notes-Only Slide",
            bullets=[], key_concepts=[], definitions=[], examples=[],
            speaker_notes="Explain the concept slowly\nUse the whiteboard\nAsk students to recall",
        )
        kit = _make_kit([slide_obj])
        prs = _run_pptx(kit, _make_course())
        content_slides = [prs.slides[i] for i in _slide_index_of_content_slides(prs)]
        all_texts = _body_texts(content_slides[0])
        combined = " ".join(all_texts)
        assert "Explain the concept slowly" in combined

    def test_alt_field_name_points_hoisted(self):
        """Content dict with 'points' (Groq alias) instead of 'bullets' is rendered."""
        slide_obj = _make_slide(
            1, "Points Alias",
            alt_field={"points": ["Point one about sorting", "Point two about searching"]},
        )
        kit = _make_kit([slide_obj])
        prs = _run_pptx(kit, _make_course())
        content_slides = [prs.slides[i] for i in _slide_index_of_content_slides(prs)]
        all_texts = _body_texts(content_slides[0])
        combined = " ".join(all_texts)
        assert "Point one about sorting" in combined

    def test_fully_empty_slide_has_placeholder_message(self):
        """A slide with absolutely no content gets a 'please edit' message — body never blank."""
        slide_obj = _make_slide(
            1, "Empty Slide",
            bullets=[], key_concepts=[], definitions=[], examples=[],
            speaker_notes=None, classroom_activity=None,
            student_summary=None, teaching_notes=None,
        )
        kit = _make_kit([slide_obj])
        prs = _run_pptx(kit, _make_course())
        content_slides = [prs.slides[i] for i in _slide_index_of_content_slides(prs)]
        all_texts = _body_texts(content_slides[0])
        # At minimum the placeholder message must be there
        combined = " ".join(all_texts)
        assert len(combined.strip()) > 10, "Body must not be blank even for fully-empty slides"

    def test_no_placeholder_xml_in_content_slides(self):
        """No slide should have unset placeholder XML (<p:ph>) in its shapes."""
        from pptx.util import Inches  # noqa: F401
        from lxml import etree

        slide_obj = _make_slide(1, "Normal Slide", bullets=["Test bullet"])
        kit  = _make_kit([slide_obj])
        prs  = _run_pptx(kit, _make_course())

        PH_TAG = "{http://schemas.openxmlformats.org/presentationml/2006/main}ph"
        for i, slide in enumerate(prs.slides):
            for shape in slide.shapes:
                sp_elem = shape._element
                ph_elems = sp_elem.findall(f".//{PH_TAG}")
                assert not ph_elems, (
                    f"Slide {i+1} shape '{shape.name}' still has placeholder XML — "
                    f"_clear_placeholders did not run or did not remove all placeholders."
                )


class TestPptxMultipleSlides:
    def test_all_content_slides_have_body_text(self):
        """For a 5-slide kit, every content slide (index >= 2) must have body text."""
        slides = [
            _make_slide(1, "Introduction to Python",
                        bullets=["Python is interpreted", "Cross-platform", "Beginner friendly"]),
            _make_slide(2, "Variables and Data Types",
                        key_concepts=["int", "str", "list", "dict"],
                        definitions=["Variable: a named memory location"]),
            _make_slide(3, "Control Flow",
                        examples=["if/elif/else for branching", "for loops for iteration"]),
            _make_slide(4, "Functions",
                        bullets=[], speaker_notes="Cover def keyword\nReturn values\nDefault args"),
            _make_slide(5, "Summary",
                        student_summary="Python basics covered: variables, control flow, functions"),
        ]
        kit = _make_kit(slides)
        prs = _run_pptx(kit, _make_course())

        content_indices = _slide_index_of_content_slides(prs)
        assert len(content_indices) >= 5, "All 5 content slides must be generated"

        for idx in content_indices:
            slide  = prs.slides[idx]
            texts  = _body_texts(slide)
            combined = " ".join(texts).strip()
            assert len(combined) > 5, (
                f"Slide at PPTX index {idx} has blank body — "
                f"shapes: {[s.name for s in slide.shapes]}"
            )

    def test_slide_title_appears_in_header(self):
        """The slide title should appear in header text (not just body)."""
        slide_obj = _make_slide(1, "My Unique Slide Title ABC", bullets=["Some content"])
        kit = _make_kit([slide_obj])
        prs = _run_pptx(kit, _make_course())
        content_slides = [prs.slides[i] for i in _slide_index_of_content_slides(prs)]
        all_texts = _body_texts(content_slides[0])
        combined = " ".join(all_texts)
        assert "My Unique Slide Title ABC" in combined
