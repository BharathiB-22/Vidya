"""
Unit tests for H-36A Phase 1 — rich PPTX teaching slide export.

Pure python-pptx verification: no DB, no Celery, no MinIO.
Tests run synchronously; asyncio_mode=auto does not affect them.
"""
from __future__ import annotations

import enum
import io
import types

import pytest

from app.workers.heavy.course_kit_export import _fill_slide_body, _generate_pptx


# ---------------------------------------------------------------------------
# Minimal mock helpers
# ---------------------------------------------------------------------------

def _course(code: str = "CS301", title: str = "Data Structures"):
    return types.SimpleNamespace(code=code, title=title)


class _FakeBloom(str, enum.Enum):
    APPLY = "APPLY"
    REMEMBER = "REMEMBER"


class _FakeStatus(str, enum.Enum):
    PUBLISHED = "PUBLISHED"


class _FakeComplexity(str, enum.Enum):
    UG = "UG"


def _slide(
    slide_number: int = 1,
    title: str = "Slide Title",
    content: dict | None = None,
    speaker_notes: str | None = None,
    bloom_level=None,
    co_reference: str | None = None,
):
    return types.SimpleNamespace(
        slide_number=slide_number,
        title=title,
        content=content if content is not None else {},
        speaker_notes=speaker_notes,
        bloom_level=bloom_level,
        co_reference=co_reference,
    )


def _kit(slides=None, teaching_plan=None, resources=None, **kw):
    return types.SimpleNamespace(
        unit_number=kw.get("unit_number", 1),
        version=kw.get("version", 1),
        status=_FakeStatus.PUBLISHED,
        complexity_level=_FakeComplexity.UG,
        slides=slides or [],
        quizlets=[],
        assignments=[],
        teaching_plan=teaching_plan or [],
        resources=resources or [],
        ai_model=None,
        published_at=None,
    )


def _render(kit, course=None, is_dean: bool = False):
    """Generate PPTX and return the parsed Presentation."""
    from pptx import Presentation

    buf = io.BytesIO()
    _generate_pptx(buf, kit, course or _course(), is_dean=is_dean)
    buf.seek(0)
    return Presentation(buf)


def _all_text(prs, slide_index: int) -> str:
    """Concatenated text from all shapes on the slide (blank-layout compatible)."""
    parts = []
    for shape in prs.slides[slide_index].shapes:
        if shape.has_text_frame:
            parts.append(shape.text_frame.text)
    return '\n'.join(parts)


# Keep old name as alias so existing tests don't need renaming.
def _body_text(prs, slide_index: int) -> str:
    return _all_text(prs, slide_index)


def _notes_text(prs, slide_index: int) -> str:
    return prs.slides[slide_index].notes_slide.notes_text_frame.text


# ---------------------------------------------------------------------------
# _fill_slide_body — direct helper tests
# ---------------------------------------------------------------------------

def test_fill_body_empty_content_no_error():
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    sl = prs.slides.add_slide(prs.slide_layouts[1])
    tf = sl.placeholders[1].text_frame
    _fill_slide_body(tf, {})  # must not raise


def test_fill_body_bullets():
    from pptx import Presentation

    prs = Presentation()
    sl = prs.slides.add_slide(prs.slide_layouts[1])
    tf = sl.placeholders[1].text_frame
    _fill_slide_body(tf, {"bullets": ["Alpha", "Beta", "Gamma"]})
    text = tf.text
    assert "Alpha" in text
    assert "Beta" in text
    assert "Gamma" in text


def test_fill_body_key_concepts():
    from pptx import Presentation

    prs = Presentation()
    sl = prs.slides.add_slide(prs.slide_layouts[1])
    tf = sl.placeholders[1].text_frame
    _fill_slide_body(tf, {"key_concepts": ["Heap", "BST"]})
    text = tf.text
    assert "Heap" in text
    assert "BST" in text
    assert "Key Concepts" in text


def test_fill_body_definitions():
    from pptx import Presentation

    prs = Presentation()
    sl = prs.slides.add_slide(prs.slide_layouts[1])
    tf = sl.placeholders[1].text_frame
    _fill_slide_body(tf, {"definitions": ["A stack is a LIFO structure."]})
    text = tf.text
    assert "stack" in text
    assert "Definitions" in text


def test_fill_body_examples():
    from pptx import Presentation

    prs = Presentation()
    sl = prs.slides.add_slide(prs.slide_layouts[1])
    tf = sl.placeholders[1].text_frame
    _fill_slide_body(tf, {"examples": ["push 1, push 2, pop → 2"]})
    text = tf.text
    assert "push 1" in text
    assert "Examples" in text


def test_fill_body_code_snippet():
    from pptx import Presentation

    prs = Presentation()
    sl = prs.slides.add_slide(prs.slide_layouts[1])
    tf = sl.placeholders[1].text_frame
    _fill_slide_body(tf, {"code_snippet": "def push(stack, item):\n    stack.append(item)"})
    text = tf.text
    assert "push" in text
    assert "Code" in text


def test_fill_body_code_snippet_capped_at_20_lines():
    from pptx import Presentation

    prs = Presentation()
    sl = prs.slides.add_slide(prs.slide_layouts[1])
    tf = sl.placeholders[1].text_frame
    long_code = "\n".join(f"line_{i}" for i in range(30))
    _fill_slide_body(tf, {"code_snippet": long_code})
    text = tf.text
    assert "line_19" in text      # last allowed line (index 19)
    assert "line_20" not in text  # first cut line


def test_fill_body_diagram_prompt():
    from pptx import Presentation

    prs = Presentation()
    sl = prs.slides.add_slide(prs.slide_layouts[1])
    tf = sl.placeholders[1].text_frame
    _fill_slide_body(tf, {"diagram_prompt": "Draw a binary search tree"})
    text = tf.text
    assert "binary search tree" in text
    assert "Diagram" in text


def test_fill_body_classroom_activity():
    from pptx import Presentation

    prs = Presentation()
    sl = prs.slides.add_slide(prs.slide_layouts[1])
    tf = sl.placeholders[1].text_frame
    _fill_slide_body(tf, {"classroom_activity": "Trace merge sort on whiteboard"})
    text = tf.text
    assert "merge sort" in text
    assert "Activity" in text


def test_fill_body_student_summary():
    from pptx import Presentation

    prs = Presentation()
    sl = prs.slides.add_slide(prs.slide_layouts[1])
    tf = sl.placeholders[1].text_frame
    _fill_slide_body(tf, {"student_summary": "A queue is FIFO, a stack is LIFO."})
    text = tf.text
    assert "FIFO" in text
    assert "Summary" in text


def test_fill_body_teaching_notes_excluded():
    """teaching_notes must never appear in the slide body."""
    from pptx import Presentation

    prs = Presentation()
    sl = prs.slides.add_slide(prs.slide_layouts[1])
    tf = sl.placeholders[1].text_frame
    _fill_slide_body(tf, {"teaching_notes": "CONFIDENTIAL_INSTRUCTOR_TEXT"})
    assert "CONFIDENTIAL_INSTRUCTOR_TEXT" not in tf.text


# ---------------------------------------------------------------------------
# _generate_pptx — integration-level tests (no DB/Celery)
# ---------------------------------------------------------------------------

def test_empty_kit_produces_cover_objectives_summary():
    # New design: cover + objectives + summary = 3 slides minimum
    prs = _render(_kit(slides=[]))
    assert len(prs.slides) == 3


def test_single_slide_produces_four_slides():
    # cover + objectives + 1 content + summary = 4 slides
    prs = _render(_kit(slides=[_slide()]))
    assert len(prs.slides) == 4


def test_slide_type_badge_on_content_slide():
    # Badge is a separate textbox (not title prefix) on blank layout
    s = _slide(title="Variable Types", content={"slide_type": "CODE"})
    prs = _render(_kit(slides=[s]))
    # Content slide is at index 2 (0=cover, 1=objectives, 2=content)
    slide_text = _all_text(prs, 2)
    assert "CODE" in slide_text
    assert "Variable Types" in slide_text


def test_no_badge_when_slide_type_absent():
    s = _slide(title="Intro", content={})
    prs = _render(_kit(slides=[s]))
    # Content slide at index 2 should have no bracket-enclosed badge
    slide_text = _all_text(prs, 2)
    assert "Intro" in slide_text
    assert "[" not in slide_text  # no type badge means no brackets


def test_bullets_in_body():
    s = _slide(content={"bullets": ["Point A", "Point B"]})
    prs = _render(_kit(slides=[s]))
    body = _body_text(prs, 2)  # content slide is index 2
    assert "Point A" in body
    assert "Point B" in body


def test_teaching_notes_in_notes_pane_for_faculty():
    s = _slide(content={"teaching_notes": "Emphasise the stack frame concept."})
    prs = _render(_kit(slides=[s]), is_dean=False)
    notes = _notes_text(prs, 2)  # content slide is index 2
    assert "TEACHING NOTES" in notes
    assert "stack frame" in notes


def test_teaching_notes_hidden_from_dean():
    s = _slide(content={"teaching_notes": "SECRET_FACULTY_NOTE"})
    prs = _render(_kit(slides=[s]), is_dean=True)
    notes = _notes_text(prs, 2)
    assert "SECRET_FACULTY_NOTE" not in notes


def test_teaching_notes_never_in_slide_body_faculty():
    s = _slide(content={"teaching_notes": "CONFIDENTIAL_BODY_CHECK"})
    prs = _render(_kit(slides=[s]), is_dean=False)
    body = _body_text(prs, 2)
    assert "CONFIDENTIAL_BODY_CHECK" not in body


def test_speaker_notes_in_notes_pane_for_faculty():
    s = _slide(content={}, speaker_notes="Say this slowly to the class.")
    prs = _render(_kit(slides=[s]), is_dean=False)
    notes = _notes_text(prs, 2)
    assert "SPEAKER NOTES" in notes
    assert "Say this slowly" in notes


def test_speaker_notes_hidden_from_dean():
    s = _slide(content={}, speaker_notes="HIDDEN_SPEAKER_NOTE")
    prs = _render(_kit(slides=[s]), is_dean=True)
    notes = _notes_text(prs, 2)
    assert "HIDDEN_SPEAKER_NOTE" not in notes


def test_bloom_and_co_always_present_even_for_dean():
    s = _slide(content={}, bloom_level=_FakeBloom.APPLY, co_reference="CO2")
    for is_dean in (True, False):
        prs = _render(_kit(slides=[s]), is_dean=is_dean)
        notes = _notes_text(prs, 2)
        assert "Bloom: APPLY" in notes
        assert "CO: CO2" in notes


def test_full_rich_slide_is_valid_pptx():
    """A slide with every H-34 field produces a parseable PPTX."""
    content = {
        "slide_type": "CONCEPT",
        "bullets": ["Bullet 1", "Bullet 2"],
        "key_concepts": ["Heap", "BST"],
        "definitions": ["A heap is a complete binary tree."],
        "examples": ["Insert 5 into max-heap: bubble up to root"],
        "code_snippet": "def heapify(arr, n, i):\n    largest = i",
        "diagram_prompt": "Draw a max-heap with 5 elements",
        "classroom_activity": "Build a min-heap on the whiteboard",
        "student_summary": "Heaps support O(log n) insert and extract-max.",
        "teaching_notes": "Note for instructor only.",
    }
    s = _slide(
        title="Heap Data Structure",
        content=content,
        speaker_notes="Speak slowly here.",
        bloom_level=_FakeBloom.APPLY,
        co_reference="CO3",
    )
    prs = _render(_kit(slides=[s]), is_dean=False)
    # cover + objectives + 1 content + summary = 4 slides
    assert len(prs.slides) == 4

    body = _body_text(prs, 2)  # content slide at index 2
    assert "Bullet 1" in body
    assert "Heap" in body
    assert "heap is a complete" in body
    assert "bubble up" in body
    assert "min-heap" in body
    assert "FIFO" not in body  # sanity: no contamination

    notes = _notes_text(prs, 2)
    assert "TEACHING NOTES" in notes
    assert "SPEAKER NOTES" in notes
    assert "Bloom: APPLY" in notes
    assert "CO: CO3" in notes
    assert "Note for instructor only" in notes
    assert "Note for instructor only" not in body  # confirmed not in body


def test_backward_compat_old_kit_only_bullets():
    """Old kits with only bullets field work without error."""
    s = _slide(content={"bullets": ["Old bullet 1"]})
    prs = _render(_kit(slides=[s]))
    assert "Old bullet 1" in _body_text(prs, 2)


def test_backward_compat_none_content():
    """Slides where content is None/empty dict render cleanly."""
    s = _slide(content=None)
    prs = _render(_kit(slides=[s]))
    assert len(prs.slides) == 4  # cover + objectives + content + summary — no crash


def test_multiple_slides_all_render():
    slides = [
        _slide(slide_number=1, title="Slide 1", content={"bullets": ["B1"]}),
        _slide(slide_number=2, title="Slide 2", content={"definitions": ["D2"]}),
        _slide(slide_number=3, title="Slide 3", content={"slide_type": "ACTIVITY",
                                                          "classroom_activity": "Group work"}),
    ]
    prs = _render(_kit(slides=slides))
    # cover(0) + objectives(1) + 3 content(2,3,4) + summary(5) = 6 slides
    assert len(prs.slides) == 6
    assert "B1" in _body_text(prs, 2)
    assert "D2" in _body_text(prs, 3)
    assert "Group work" in _body_text(prs, 4)
    # ACTIVITY badge is a textbox on blank layout, not in the title placeholder
    assert "ACTIVITY" in _all_text(prs, 4)
